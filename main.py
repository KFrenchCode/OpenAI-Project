from pprint import pp
from string import capwords
import tkinter as tk
from tkinter import ttk
from tkinter import *
from tkinter import Tk, Label, PhotoImage
from time import sleep
from tkinter import Canvas as tkCanvas, filedialog
from tkinter import messagebox
import tkinter
from turtle import width
from xml.sax.handler import property_xml_string
from click import wrap_text
from transformers import pipeline, Pipeline, PreTrainedTokenizerFast, BartTokenizer,BartForConditionalGeneration
import requests
from bs4 import BeautifulSoup, Comment
import fitz
from openai import OpenAI
from dotenv import load_dotenv
import tiktoken
from docx import Document
from reportlab.pdfgen.canvas import Canvas as pdfCanvas
from reportlab.lib.pagesizes import LETTER
import textwrap
from textwrap import wrap
import datetime
from threading import Thread
import pkg_resources
import threading
import queue
from pptx import Presentation
from tkinter import messagebox
import llama_index
import chromadb
from typing import Optional, List
from sentence_transformers import SentenceTransformer
from llama_index.embeddings.openai import OpenAIEmbedding
from llama_index.embeddings.huggingface import HuggingFaceEmbedding
from llama_index.core import VectorStoreIndex, SimpleDirectoryReader, Settings
import pickle
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.cluster import KMeans
from llama_index.embeddings.huggingface import HuggingFaceEmbedding
from llama_index.core import Settings
import os

import openai

#hard code exporting API key here.    

# Step 1: Hard code your actual API key (make sure it's a string)
api_key = "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"
os.environ["MY_KEY"] = api_key
# Step 3: Verify that the API key is set (optional)
# print(f"API_KEY is set to: {os.getenv('MY_KEY')}")     
# Disable parallelism for huggingface tokenizers
os.environ["TOKENIZERS_PARALLELISM"] = "false"
# Access the API key from the environment variable
api_key = os.environ.get('MY_KEY')

# Initialize the OpenAI API client
openai.api_key = api_key

Settings.embed_model = HuggingFaceEmbedding(model_name="BAAI/bge-small-en-v1.5")





def rag(query: Optional[str] = None):
    # Load documents
    documents = SimpleDirectoryReader("data").load_data()
    
    # Initialize the index with documents
    index = VectorStoreIndex.from_documents(documents, embed_model='local')
    
    # Create a query engine from the index
    query_engine = index.as_query_engine()
    
    # Retrieve relevant documents based on the query
    retrieved_documents = query_engine.query(query)
    
    # Generate embeddings for the retrieved documents
    def generate_embeddings(documents):
        model = SentenceTransformer('paraphrase-MiniLM-L6-v2')
        embeddings = model.encode(documents)
        return embeddings
    
    # Generate a response using the retrieved documents
    def generate_response(retrieved_documents):
        response = " ".join(retrieved_documents)
        return response
    
    # Apply the analytical framework to the retrieved documents
    def apply_analytical_framework(documents):
        # Vectorize the documents
        vectorizer = TfidfVectorizer(stop_words='english')
        X = vectorizer.fit_transform(documents)
        
        # Apply KMeans clustering to group documents into categories
        kmeans = KMeans(n_clusters=5, random_state=0).fit(X)
        
        # Assign categories to documents
        categories = kmeans.labels_
        categorized_documents = {i: [] for i in range(5)}
        for doc, category in zip(documents, categories):
            categorized_documents[category].append(doc)
        
        return categorized_documents
    
    # Apply the analytical framework to the retrieved documents
    categorized_documents = apply_analytical_framework(retrieved_documents)
    
    # Generate a response using the categorized documents
    response = generate_response(categorized_documents)
    
    return response

# Example usage
documents = SimpleDirectoryReader("Data").load_data()
index = VectorStoreIndex.from_documents(documents)
query_engine = index.as_query_engine()
response = query_engine.query("How do these documents shape the analysis of the given source?")

print(response)

# Storing the Vector Index
# Serialize the index for storage
with open('index.pkl', 'wb') as f:
    pickle.dump(index, f)

# Loading the Vector Index
# Deserialize the index from storage
with open('index.pkl', 'rb') as f:
    index_loaded = pickle.load(f)

# Use the loaded index
query_engine_loaded = index_loaded.as_query_engine()
response_loaded = query_engine_loaded.query("How do these documents shape the analysis of the given source?")

print(response_loaded)



#BG (blue = #6EB3F4), + image in python-first-steps, bg.  

tokenizer = BartTokenizer.from_pretrained('facebook/bart-large-cnn')
model = BartForConditionalGeneration.from_pretrained('facebook/bart-large-cnn')


load_dotenv()


class ReportGeneratorApp:
    def __init__(self, root):

        self.window = root
        # Initialize variables
        self.root: tk.Tk = root
        self.bg_image = PhotoImage(file="/Users/kendrafrench/Dev/python-first-steps/BG/bgfinal_nb.png")
        # create a label 
        bg_label = Label(root, image = self.bg_image)
        bg_label.place(x=0,y=0, relwidth=1,relheight=1)

        self.source_list: list[dict[str, str]] = []
        self.source_widgets: list[tk.Widget] = []
        self.title: str = "CCJ2-APLE Virtual Analyst"

        
        
       

        # Initialize source type variable
        self.source_type_var = tk.StringVar()
        

        # Initialize widgets
        self.root.title(self.title)
        self.create_initial_widgets()

        self.current_source_location: str = ""

        # # Initialize Model LEGACY, SWITCHED TO OPENAI
        # self.summarizer: Pipeline = pipeline("summarization", model="facebook/bart-large-cnn")
        
        # Initialize OpenAI
        self.client = OpenAI(
            api_key=os.getenv("MY_KEY"),
        )

        # Initialize encoder
        self.encoder = tiktoken.get_encoding("cl100k_base")

        self.export_document = Document()

        # Initialize Summaries
        self.summaries: list[dict[str, str]] = []

        self.summary_object_directions = ""


        # self.root = tk.Tk() 
        self.export_pdf = pdfCanvas("ReportTemplate-pdf", pagesize=LETTER)


    def save_summaries_to_pdf(self):
        # self.save_pdf_text_variable.set("Saving...")
        # self.save_pdf_text_variable["state"]=tk.DISABLED
        # need to add coding to the button itself

        print("Report saved to PDF")
        canvas = pdfCanvas("ReportTemplate.pdf", pagesize=LETTER)

        # heading formatting
        canvas.setFont("Helvetica", 18)
        canvas.drawString(72,740, "[Intelligence Note or Reporting Highlights]")
        canvas.setFont("Helvetica", 12)
            
        # formatting contents
        y = 700
        max_chars_per_line = 90 # Adjust this value as needed
        for summary in self.summaries:
            classification = summary["source_classification"]
            country = summary["source_country"]
            title = summary["source_title"]
            summary_text = summary["source_summary"]
            citation = summary["source_citation"]

            canvas.drawString(72, y, f"Classification: {classification}")
            y -= 20
            canvas.drawString(72, y, f"Country: {country}")
            y -= 20
            canvas.drawString(72, y, f"Title: {title}")
            y -= 40
            # canvas.drawString(72, y, f"{summary_text}")
            # split the text into lines
            summary_lines = wrap(summary_text, width=max_chars_per_line)
            for i, line in enumerate(summary_lines):
                canvas.drawString(72, y - i*20, line)
            y -= len(summary_lines)*20 + 60
            y -= 0

            canvas.drawString(72, y, f"Analyst Comment:")
            y = -20
            # split the text into lines
            citation_lines = wrap(citation, width=max_chars_per_line)
            for i, line in enumerate(citation_lines):
                canvas.drawString(72, y - i*20, line)
            y -= len(citation_lines)*20 + 60
            y -= 0


            # self.save_pdf_text_variable.set("PDF Saved")
            # self.save_pdf_text_variable["state"]=tk.NORMAL

            canvas.save()

            # Call loading_page from the main thread
            # threading.Timer(0, loading_page, args=(self, self.save_summaries_to_pdf)).start()

            

    def save_summaries_to_docx(self):
        self.summarize_button_text_variable.set("Thinking...")
        self.summarize_button["state"]=tk.DISABLED

        print("Report Saved to Word File")
        self.export_document.add_heading("[Intelligence Note or Reporting Highlights]", level=0)
        for summary in self.summaries:
            classification = self.export_document.add_paragraph()
            classification.add_run(summary["source_classification"]).bold = True
            # add caps to classification
            country = self.export_document.add_paragraph()
            country.add_run(summary["source_country"]).bold = True
            title = self.export_document.add_paragraph()
            title.add_run(summary["source_title"]).bold = True

            self.export_document.add_paragraph()

                
            self.export_document.add_paragraph(summary["source_summary"])
            
            add_analystc_comment = self.export_document.add_paragraph()

            self.export_document.add_paragraph()

            add_analystc_comment.add_run("[Analyst Comment]").bold = True

            citation = self.export_document.add_paragraph()
            citation.add_run(summary["source_citation"])


            self.summarize_button_text_variable.set("Document Saved")
            self.summarize_button["state"]=tk.NORMAL

            
        # Export a document loading page
        # Call loading_page from the main thread
        # threading.Timer(0, loading_page, args=(self, self.save_summaries_to_docx)).start()

        


            self.export_document.save("ReportTemplate.docx")
            self.summarize_button_text_variable.set("Summarize Sources")
            self.summarize_button["state"]=tk.NORMAL


    #create/ Save Source Packet


    def create_source_packet(self):
        source_packet_popup = tk.Toplevel(self.window)
        source_packet_popup.title("Source Packet By Virtual Analyst Preview")
        source_packet_popup.geometry('1000x1000')

        import openai
        from docx import Document
        from docx.shared import RGBColor
        from docx.enum.text import WD_COLOR_INDEX

                    
        # Set your OpenAI API key
        openai.api_key = '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'

        def chat_with_gpt(prompt, chat_history):
            response = openai.ChatCompletion.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "user", "content": prompt},
                    {"role": "assistant", "content": chat_history}
                ]
            )

            return response['choices'][0]['message']['content']

        # print("Welcome to ChatGPT. You can start chatting by typing your messages.")
        # print("Type 'exit' to end the conversation.")

        chat_history = ""

        source = ""
        product = ""

        try:
            with open('source.txt', 'rb') as file:
                file_contents = file.read()
                # Decode the binary content with a specific encoding
                decoded_contents = file_contents.decode('utf-8')  # Replace 'utf-8' with the appropriate encoding
            source = decoded_contents
        except FileNotFoundError:
            print("File not found.")
        except UnicodeDecodeError:
            print("Error: Unable to decode the file with the specified encoding.")

        try:
            with open('product.txt', 'rb') as file:
                file_contents = file.read()
                # Decode the binary content with a specific encoding
                decoded_contents = file_contents.decode('utf-8')  # Replace 'utf-8' with the appropriate encoding
            product = decoded_contents
        except FileNotFoundError:
            print("File not found.")
        except UnicodeDecodeError:
            print("Error: Unable to decode the file with the specified encoding.")

        user_input = "This is the news source: " + source

        response = chat_with_gpt(user_input, chat_history)
        chat_history += user_input + "\n"
        chat_history += response + "\n"
        # print("ChatGPT:", response)

        user_input = "This is the analysis of the new source: " + product

        response = chat_with_gpt(user_input, chat_history)
        chat_history += user_input + "\n"
        chat_history += response + "\n"
        # print("ChatGPT:\n", response)

        user_input = """I want you to list all of the quotes from the news source that was referenced in the analysis. 
                        Do not change or remove anything to the quotes from the news source. Do not include the quotations around the quotes.
                        Do not include duplicate quotes. Do not add on anything besides the quotes. 
                        The list must be numbered. The quotes should come in the order it shows up in the text"""


        response = chat_with_gpt(user_input, chat_history)
        chat_history += user_input + "\n"
        chat_history += response + "\n"
        print("ChatGPT:\n", response)

        quotes_list = response.split("\n")
        print(quotes_list)
        print(len(quotes_list))

        for i in range(len(quotes_list)):
            quotes_list[i] = quotes_list[i][4:len(quotes_list[i])-3]

        print(quotes_list)

        def create_word_doc(text, keywords, title, author, date):
            doc = Document()
            doc.add_heading(title, 0)
            doc.add_heading(author + " " + date, 1)
            paragraph = doc.add_paragraph("")
            text_index = 0
            for keyword in keywords:
                if keyword in text:
                    index = text.find(keyword)
                    paragraph.add_run(text[text_index: index])
                    paragraph.add_run(text[index:index + len(keyword)]).font.highlight_color = WD_COLOR_INDEX.YELLOW 
                    text_index = index + len(keyword)
            paragraph.add_run(text[text_index:])
            doc.save('highlighted_doc.docx')

        create_word_doc(source, quotes_list, "Testing", "daniel tran", "12 15 2001")



        from docx import Document
        from docx.shared import RGBColor
        from docx.enum.text import WD_COLOR_INDEX


        def highlight_text(paragraph, text_to_highlight):
            for run in paragraph.runs:
                if text_to_highlight in run.text:
                    start_index = run.text.find(text_to_highlight)
                    end_index = start_index + len(text_to_highlight)
                    new_run = paragraph.add_run(run.text[start_index:end_index])
                    new_run.font.highlight_color = WD_COLOR_INDEX.YELLOW

        def create_highlighted_document(original_source, derivative):
            doc = Document()

            derivative_paragraphs = derivative.split("\n")

            for derivative_paragraph in derivative_paragraphs:
                doc.add_paragraph(derivative_paragraph)

            for paragraph in doc.paragraphs:
                for derivative_paragraph in derivative_paragraphs:
                    if derivative_paragraph in paragraph.text:
                        highlight_text(paragraph, derivative_paragraph)

            return doc

        def save_document(document, filename):
            document.save(filename)

        if __name__ == "__main__":
            original_source = """
            President Joe Biden on Thursday told Israeli Prime Minister Benjamin Netanyahu that the overall humanitarian situation in Gaza is unacceptable and warned Israel to take steps to address the crisis or face consequences, a stark statement from Israel’s staunchest ally.

        The 30-minute conversation was the two leaders’ first phone call since an Israeli strike killed seven aid workers from the World Central Kitchen in Gaza. The killings have set off a furor inside the White House and Biden has been described as reaching a new level of frustration with Israel’s campaign.

        “President Biden emphasized that the strikes on humanitarian workers and the overall humanitarian situation are unacceptable,” the White House said in a statement shortly after the call wrapped. “He made clear that US policy with respect to Gaza will be determined by our assessment of Israel’s immediate action on these steps.”

        Biden also said Israel needed to “announce and implement a series of specific, concrete, and measurable steps to address civilian harm, humanitarian suffering, and the safety of aid workers.”

        Speaking in Brussels after the call, Secretary of State Antony Blinken made the stakes clear: “If we don’t see the changes that we need to see, there’ll be changes in our own policy,” he said. However, neither Blinken nor National Security Council spokesman John Kirby, speaking later at the White House press briefing, detailed what those potential policy changes could be.

        Kirby said the US “would hope to see some announcements of changes here in the coming hours and days – and I’ll leave it at that.”

        The call amounts to perhaps the most serious sign of Biden’s frustration with Israel’s campaign in Gaza, which was launched in the wake of the October 7 attacks by Hamas. In the time since, more than 32,000 people have been killed in Gaza, according to the Palestinian enclave’s health ministry, and grisly reports of civilian deaths are made every day.

        The war has become one of Biden’s chief domestic political problems ahead of November’s election as key parts of his voter coalition have been outraged by the president’s support for Israel’s war. Protests have popped up at nearly every public event the president has held outside the White House in recent months and he was confronted in an intimate meeting with Muslim leaders earlier this week with stark opposition to the US’ policy toward the war.

        That political conundrum made Thursday’s call and the statements from Blinken and Kirby particularly notable, as both men indicated possible change in the US’ posture toward the war may come if Israel keeps up its current practices.

        Blinken said Israel should not stoop to the level of Hamas in its response to the group’s October 7 attacks.

        “Israel is not Hamas,” Blinken said. “Israel is a democracy, Hamas a terrorist organization. And democracies place the highest value on human life - every human life.”

        Biden was shaken by the strike on the World Central Kitchen workers, Kirby said.

        Israel has acknowledged responsibility for the strikes but said the convoy was not targeted and the workers’ deaths were not intentional. The country continues to investigate the circumstances surrounding the killings. Blinken acknowledged the World Central Kitchen strike was not the first time Israel has killed aid workers in the conflict.

        “It must be the last,” he said.

        But even at the same time that the administration was alluding to potential changes in policy if Israel doesn’t stop killing civilians, the US was sending more deadly weapons to its ally.

        Biden is set to green light an $18 billion sale of fighter jets from the United States and Israel and the administration recently authorized the transfer to Israel of over 1,000 500-pound bombs and over 1,000 small-diameter bombs, according to three people familiar with the matter.

        Kirby defended the arms sale and transfers as the product of yearslong processes.

        “With the exception of the immediate two months after the attack, we haven’t really sent emergency aid and military assistance to Israel,” Kirby said. “What you’re seeing here is the result of a process of foreign military sales to Israel that takes years.”

        Earlier on Thursday, one of Biden’s closest allies on Capitol Hill – Democratic Sen. Chris Coons – told CNN that for the first time he is open to applying conditions to aid for Israel after the strike on the World Central Kitchen workers.

        “I think we’re at that point,” he told CNN’s Sara Sidner. “I think we’re at the point where President Biden has said – and I have said and others have said – if Benjamin Netanyahu, prime minister, were to order the IDF into Rafah at scale, they were to drop thousand pound bombs and send in a battalion to go after Hamas and make no provision for civilians or for humanitarian aid, that I would vote to condition aid to Israel.”

        The president has called for a temporary ceasefire that includes the release of hostages held by Hamas, and has repeatedly said he does not want Israel to undertake a ground invasion of Rafah in southern Gaza, where scores of displaced Gazans have sought shelter. He again called for a temporary ceasefire in the call with Netanyahu on Thursday.

        “He underscored that an immediate ceasefire is essential to stabilize and improve the humanitarian situation and protect innocent civilians, and he urged the prime minister to empower his negotiators to conclude a deal without delay to bring the hostages home,” the White House said in its statement following the call.

        But, Biden has so far stopped short of calling for a permanent ceasefire. His hesitance to do so is increasingly out of step with the actions and public statements of other world leaders, including US allies.

        This is a breaking story and will be updated

        CNN’s Jennifer Hansler and Donald Judd contributed to this report.
            """

            derivative = """
            The recent phone call between President Joe Biden and Israeli Prime Minister Benjamin Netanyahu marks a significant shift in the relationship dynamics between the two allies. Biden's assertion of the unacceptability of the strikes on humanitarian workers in Gaza reflects a heightened level of frustration within the Biden administration regarding Israel's military campaign.

        The White House's statement following the call underscores the seriousness of the situation, indicating that US policy towards Gaza will be contingent upon Israel's immediate actions to address civilian harm and humanitarian suffering. This stance suggests that the US may be considering changes in its policy towards Israel if these concerns are not adequately addressed.

        Secretary of State Antony Blinken's remarks further emphasize the gravity of the situation, indicating that failure to see the necessary changes may result in alterations to US policy. Blinken's assertion that Israel, as a democracy, should uphold the highest value of human life underscores the moral dimension of the issue.

        However, despite the rhetoric and expressions of concern, the US's continued arms sales to Israel raise questions about the consistency of its position. The sale of fighter jets and bombs to Israel, despite the ongoing humanitarian crisis, highlights the complexity of US-Israel relations and the challenges of balancing strategic interests with humanitarian concerns.

        Senator Chris Coons' openness to applying conditions to aid for Israel reflects a growing sentiment within the Democratic Party to reassess the unconditional support traditionally extended to Israel. This shift in attitude suggests a potential divergence between the Biden administration and some members of Congress regarding the approach to US-Israel relations.

        Overall, while President Biden's call for a temporary ceasefire and emphasis on humanitarian concerns demonstrates a willingness to engage diplomatically, the lack of a clear stance on a permanent ceasefire raises questions about the administration's long-term strategy and alignment with international consensus on resolving the conflict.

        Please note that this analysis is based on the information available up to January 2022, and subsequent developments may have occurred since then.
            """

        highlighted_document = create_highlighted_document(original_source, derivative)
        save_document(highlighted_document, "highlighted_news.docx")

        scroll_bar = Scrollbar(root)
        scroll_bar.pack( side = RIGHT, 
                            fill = Y)
        mylist = Listbox(root,
                            yscrollcommand= scroll_bar.set)
        for line in range(1,26):
            mylist.insert(END, "content" + str(line))

        mylist.pack( side = LEFT, fill = BOTH )

        scroll_bar.config( command = mylist.yview)

        include = tk.Label(source_packet_popup, text = derivative, wraplength=500, width=150, justify=tk.LEFT)
        buttonz = tk.Button(source_packet_popup, text = "Done", command=source_packet_popup.destroy)
        include.pack()
        buttonz.pack()
    
    def save_source_packet(self):
            self.save_ats_to_docx_button.set("Thinking...")
            self.save_source_packet_button["state"]=tk

            print("Source Packet Saved to Word File")
            self.export_document.add_heading("[Source Packet for Report Generated]", level=0)
            for summary in self.summaries:
                # classification = self.export_document.add_paragraph()
                # title.add_run(summary["source_title"]).bold = True

                self.export_document.add_paragraph()

                    
                self.export_document.add_paragraph(summary["create_source_packet"])
            

                self.save_source_packet_button.set("Source Saved")
                self.summarize_button["state"]=tk.NORMAL

                
            # Export a document loading page
            # Call loading_page from the main thread
            # threading.Timer(0, loading_page, args=(self, self.save_summaries_to_docx)).start()

                self.export_document.save("SourcePacket.docx")
                self.save_source_packet_button.set("Summarize Sources")
                self.save_source_packet_button["state"]=tk.NORMAL

    def Save_Generate_Briefing(self):
        self.Save_Generate_Briefing["text"] = "Thinking..."
        self.Save_Generate_Briefing["state"] = tk.DISABLED


        print("Briefing Saved to Word File")
        self.export_document.add_heading("[Briefing Script Generated]", level=0)
        for summary in self.summaries:
            # classification = self.export_document.add_paragraph()
            # title.add_run(summary["source_title"]).bold = True

            # if self.summary_object_formatting is None
            #     return
            make_brief = self.client.chat.completions.create(
                messages=[
                    {
                    "role": "system",
                    "content": """Write a Information Briefing based on the pre-fed text

                    The information briefing focuses on the clear and useful communication of facts and information. 

                    The information will always have a background: how current it is, how reliable it is, and why you are communicating it—including how it meets the needs and expectations of those you are briefing. 

                    Examples of information appropriate for an information briefing include: 
                    high-priority or time-sensitive information requiring immediate attention 
                    complex information, including complicated plans, procedures, systems, situations, statistics
                    charts requiring detailed explanation
                    controversial information requiring elaboration and explanation 

                    The information briefing’s intent is to inform the audience as efficiently as possible, so you will use a standard format. 

                    The Format & Steps for an Information Briefing 
                    Greeting & Address: the person(s) you are briefing. Identify yourself and your organization. 

                    Type and Classification: Explain the briefing’s classification. “This is a SECRET information briefing,” or “This is an UNCLASSIFIED information briefing.” If unknown, state UNCLASSIFIED unless specified otherwise.

                    Purpose and Scope: State the purpose first: Explain your briefing’s purpose and scope. “The purpose of this briefing is to bring you up to date on the company perimeter defense plan,” or “I will cover the operations plan for FTX Team Spirit.” 

                    Outline or Procedure: Briefly summarize the main point, key supporting points, and your general approach. Explain any special procedures (demonstrations, displays, or tours). Example: “During my briefing, I’ll discuss the four phases of our plan. I’ll refer to maps of our sector, and then my assistant will bring out a sand table to show you the expected sequence of the movement to contact.” 

                    Body: Arrange the main points in a logical sequence. Use visual aids correctly to emphasize your main and supporting ideas. Plan effective transitions from one main point to the next. Be prepared to answer questions at any time.

                    Closing: Ask for questions. Briefly restate your main ideas and make a concluding statement.

                    Be able to at command shorten or lengthen the brief to five, ten, or 15 minutes. The default brief will be five minutes, but maintain the quality of the argument, regardless of brief time. 

                        
                    """
                    },
                    {
                        "role": "user",
                        "content": self.summary_object_formatting.choices[0].message.content


                    }
                ],

                model="gpt-3.5-turbo-16k"
            )

            make_brief = make_brief.choices[0].message.content

            self.export_document.add_paragraph()

                
            self.export_document.add_paragraph(summary["briefing"])
          

            self.Save_Generate_Briefing.set("Briefing Script Saved", state = tk.NORMAL)

            
        # Export a document loading page
        # Call loading_page from the main thread
        # threading.Timer(0, loading_page, args=(self, self.save_summaries_to_docx)).start()

        


            self.export_document.save("TalkingPoints.docx")
            self.save_pts.set("Summarize Sources")
            self.save_pts["state"]=tk.NORMAL

    def create_initial_widgets(self):
        # Create title
        title_label = tk.Label(self.root, text=self.title, font=("Arial", 25))
        title_label.grid(column=0, row=0)

        # New Source Label
        self.source_list_label = tk.Label(self.root, text="Add Source", font=("Arial", 18))
        self.source_list_label.grid(column=0, row=1)

        # Create first source frame
        new_source_frame = tk.Frame(self.root )
        new_source_frame.grid(column=0, row=2)

        # Title of Source Label
        new_source_title_label = tk.Label(new_source_frame, text="Title of Source:", font=("Arial", 12))
        new_source_title_label.grid(column=0, row=0, pady=5)

        # Title of Source Entry
        self.new_source_title_entry = tk.Entry(new_source_frame)
        self.new_source_title_entry.grid(column=1, row=0, pady=5)

        # Country of Source Label
        new_source_country_label = tk.Label(new_source_frame, text="Country of Source:", font=("Arial", 12), width=50)
        new_source_country_label.grid(column=0, row=1, pady=5)

        # Country of Source Entry
        self.new_source_country_entry = tk.Entry(new_source_frame)
        self.new_source_country_entry.grid(column=1, row=1, pady=5)

        # Source Type Label
        source_type_label = tk.Label(new_source_frame, text="Source Type:", font=("Arial", 12))
        source_type_label.grid(column=0, row=2, pady=5)

        # Radio button for URL
        self.source_type_url = tk.Radiobutton(new_source_frame, text="URL", variable=self.source_type_var, value="url", command=self.toggle_source_input) 
        self.source_type_url.grid(column=1, row=2, pady=5)

        # Radio button for File
        self.source_type_file = tk.Radiobutton(new_source_frame, text="File", variable=self.source_type_var, value="file", command=self.toggle_source_input)
        self.source_type_file.grid(column=1, row=3, pady=5)

        # Initially, set the source type to URL
        self.source_type_var.set("url")

        # URL of Source Label
        new_source_url_label = tk.Label(new_source_frame, text="URL:", font=("Arial", 12))
        new_source_url_label.grid(column=0, row=4, pady=5)

        # # URL of Source Entry
        self.new_source_url_entry = tk.Entry(new_source_frame)
        self.new_source_url_entry.grid(column=1, row=4, pady=5)

        # # File Upload Button
        self.file_upload_button = tk.Button(new_source_frame, text="Select File", command=self.upload_file)
        # self.file_upload_button.grid(column=0, row=4, pady=5)

        # (Endnote classification) Originator; Source identifier; Date of publication; (skipping 1/17) Date of information [optional but preferred]; (Classification of title/subject) Title/Subject; (skipping 1/17) Page/paragraph or portion indicator [required when applicable]; Classification of extracted information is “X”; Overall classification is “X”; Access date.

        # URL of Source Label
        originator_label = tk.Label(new_source_frame, text="Originator:", font=("Arial", 12))
        originator_label.grid(column=0, row=5, pady=5)

        # Originator Entry
        self.originator_entry = tk.Entry(new_source_frame)
        self.originator_entry.grid(column=1, row=5, pady=5)

        # Date of Publication Label
        date_publication_label = tk.Label(new_source_frame, text="Date of Publication:", font=("Arial", 12))
        date_publication_label.grid(column=0, row=6, pady=5)

        # Date of Publication Entry
        self.date_publication_entry = tk.Entry(new_source_frame)
        self.date_publication_entry.grid(column=1, row=6, pady=5)

        classification_options = [
            "Unclassified",
            "Confidential",
            "Secret",
            "Top Secret",
        ]

        # Class. Title Label
        classification_title_label = tk.Label(new_source_frame, text="Classification of Subject/Title:", font=("Arial", 12))
        classification_title_label.grid(column=0, row=7, pady=5)

        # Class. Title Entry
        self.classification_title_var = tk.StringVar()
        self.classification_title = tk.OptionMenu(new_source_frame, self.classification_title_var, *classification_options)
        self.classification_title.grid(column=1, row=7, pady=5)

        # Class. Title Label
        portion_classification_label = tk.Label(new_source_frame, text="Classification of Portion:", font=("Arial", 12))
        portion_classification_label.grid(column=0, row=8, pady=5)

        # Class. Portion Entry
        self.portion_classification_var = tk.StringVar()
        self.portion_classification = tk.OptionMenu(new_source_frame, self.portion_classification_var, *classification_options)
        self.portion_classification.grid(column=1, row=8, pady=5)

        # Class. Title Label
        overall_classification_label = tk.Label(new_source_frame, text="Overall Classification:", font=("Arial", 12))
        overall_classification_label.grid(column=0, row=9, pady=5)

        # Class. Overall Entry
        self.overall_classification_var = tk.StringVar()
        self.overall_classification = tk.OptionMenu(new_source_frame, self.overall_classification_var, *classification_options)
        self.overall_classification.grid(column=1, row=9, pady=5)


        # Submit Button
        new_source_submit = tk.Button(new_source_frame, text="Add Source", command=self.add_new_source_command)
        new_source_submit.grid(column=0, row=10, pady=5)


        #Source List Label
        self.source_list_label = tk.Label(self.root, text="Source List", font=("Arial", 18), width=50)
        self.source_list_label.grid(column=1, row=1)

        # Source List Frame
        self.source_list_frame = tk.Frame(self.root)
        self.source_list_frame.grid(column=1, row=2, sticky="nsew")



        # Save Summaries Frame 
        self.save_source_summaries_frame = tk.Frame(self.root)
        self.save_source_summaries_frame.grid(column=0, row=4)

        # Summarize Button
        self.summarize_button_text_variable = tk.StringVar()
        self.summarize_button_text_variable.set("Generate Report")
        self.summarize_button = tk.Button(self.root, textvariable=self.summarize_button_text_variable , command=self.summarize)
        self.summarize_button.grid(column=0, row=3, pady=6)

        #Preview Report
        self.preview_report_button = tk.Button(self.save_source_summaries_frame, text="View Report" , command=self.popup)
        self.preview_report_button.grid(column=0, row=1, pady=6)

        # # Use ATS Button
        self.use_ats_button = tk.Button(self.save_source_summaries_frame, text = "View ATS", command=self.ats_popup)
        self.use_ats_button.grid(column=1, row=2, pady=10)

        # Save ATS to docx Button
        self.save_ats_to_docx_button =tk.Button(self.save_source_summaries_frame, text = "Save ATS to Word Document" , command=self.save_ats_to_docx)
        self.save_ats_to_docx_button.grid(column=1, row=3, pady=8)

         # # Intelligence Paper Button
        self.view_intel_paper = tk.Button(self.save_source_summaries_frame, text = "View Intelligence Paper", command=self.intel_paper)
        self.view_intel_paper.grid(column=1, row=4, pady=12)

         # Save Intel Paper to docx Button
        self.save_intel_paper_docx =tk.Button(self.save_source_summaries_frame, text = "Save Intelligence Paper to Word Document" , command=self.save_intel_paper)
        self.save_intel_paper_docx.grid(column=1, row=5, pady=10)

        # # New Source Label
        # self.save_source_summaries_label = tk.Label(self.save_source_summaries_frame, text="Save Source", font=("Arial", 16),bg="#6EB3F4")
        # self.save_source_summaries_label.grid(column=0, row=0)

          # # Talking Points Button
        self.view_talking_pts = tk.Button(self.save_source_summaries_frame, text = "View Talking Points", command=self.view_pts)
        self.view_talking_pts.grid(column=3, row=2, pady=12)

         # Save Talking Points to docx Button
        self.save_talking_pts =tk.Button(self.save_source_summaries_frame, text = "Save Talking Points" , command=self.save_pts)
        self.save_talking_pts.grid(column=3, row=3, pady=10)

        #Save to DOC
        
        self.save_button_text_variable = tk.StringVar()
        self.summarize_button_text_variable.set("Summarize Sources")
        self.save_source_summaries_button = tk.Button(self.save_source_summaries_frame, text="Save Report to Word Document" , command=self.save_summaries_to_docx)
        self.save_source_summaries_button.grid(column=0, row=2, pady=5)

        #Save to PDF
        self.save_source_summaries_button_pdf = tk.Button(self.save_source_summaries_frame, text="Save Report to to PDF", command=self.save_summaries_to_pdf)
        self.save_source_summaries_button_pdf.grid(column=0, row=3, pady=6)

         #Save to pptx
        
        # self.save_pptx_button_text_variable = tk.StringVar()
        # self.summarize_button_text_variable.set("Summarize Sources")
        # self.save_source_summaries_button = tk.Button(self.save_source_summaries_frame, text="Save to PowerPoint", command=self.save_summaries_to_pptx)
        # self.save_source_summaries_button.grid(column=0, row=4, pady=7)


          # # Create Source Packet Button
        self.create_soure_packet_button = tk.Button(self.save_source_summaries_frame, text = "Generate Source Packet", command=self.create_source_packet)
        self.create_soure_packet_button.grid(column=4, row=2, pady=12)

         # Save Talking Points to docx Button
        self.save_source_packet_button =tk.Button(self.save_source_summaries_frame, text = "Save Source Packet" , command=self.save_source_packet)
        self.save_source_packet_button.grid(column=4, row=3, pady=10)

         # Save/Create Briefing to docx Button
        self.create_briefing =tk.Button(self.save_source_summaries_frame, text = "Generate and Save Briefing" , command=self.Save_Generate_Briefing)
        self.create_briefing.grid(column=4, row=4, pady=12)

        # intel_options = [
        #     "Warning Intelligence",
        #     "Current Intelligence",
        #     "General Military Intelligence",
        #     "Target Intelligence",
        #     "Estimative Intelligence",

        # ]

        # # Intel Type Label
        # intel_type_label = tk.Label(self.save_source_summaries_frame, text="Intelligence Type:", font=("Arial", 14))
        # intel_type_label.grid(column=3, row=5,pady=5)


        # # Intel Type Dropdown
        # self.intel_type_var = tk.StringVar()
        # self.intel_type = tk.OptionMenu(self.save_source_summaries_frame, self.intel_type_var, *intel_options)
        # self.view_warning_intel = tk.OptionMenu(self.save_source_summaries_frame, command=self.view_warning_intel)
        # self.intel_type.grid(column=4, row=5, pady=5)




         # Disclaimer Label
        self.disclaimer_label = tk.Label(self.root, text="     * Intended for CENTCOM use only. NO classified data should be input through the system. * \n \n" 
                                         
                                                                  "    * The analytical standards and writing style used are up to date as of January 2022 but will not update until the system itself is updated.*" , font=("Arial", 14) ) 
        self.disclaimer_label.grid(column=0, row=10, pady=15)




    def toggle_source_input(self):
        source_type = self.source_type_var.get()

        # Hide all widgets initially
        self.new_source_url_entry.grid_forget()
        self.file_upload_button.grid_forget()

        # Show the relevant widget based on the source type
        if source_type == "url":
            self.new_source_url_entry.grid(column=1, row=4, pady=5, padx=5)
        elif source_type == "file":
            self.file_upload_button.grid(column=1, row=4, pady=5, padx=5)

    def create_file_upload_button(self):
        # Create a file upload button
        self.file_upload_button = tk.Button(self.root, text="Select File" ,command=self.upload_file)
        self.file_upload_button.grid(row=4, pady=5, padx=5)

    def upload_file(self):
        # Open a file dialog to select a file
        self.current_source_location = filedialog.askopenfilename(filetypes=[("PDFs", ".pdf"), ("Word documents", ".docx")])


    def reset_source_inputs(self) -> None:
        self.new_source_title_entry.delete(0, 'end')
        self.new_source_url_entry.delete(0, 'end')

    def delete_source_command(self, index: int) -> None:
        # Delete the source at the specified index
        del self.source_list[index]
        
        # Update the GUI to reflect the changes
        self.update_source_list_gui()

    def update_source_list_gui(self) -> None:
        # Clear all widgets in source_list_frame
        for widget in self.source_list_frame.winfo_children():
            widget.destroy()

        # Re-populate source_list_frame with updated source_list
        for i, source in enumerate(self.source_list):
            source_title = source["source_title"]
            
            # Source List Item Title
            source_list_title = tk.Label(self.source_list_frame, text=source_title, font=("Arial", 12))
            source_list_title.grid(column=0, row=i, sticky="w")

            # Delete Button
            delete_button = tk.Button(self.source_list_frame, text="Delete", command=lambda i=i: self.delete_source_command(i))
            delete_button.grid(column=1, row=i, padx=5, sticky="e")

    def add_new_source_command(self) -> None:
        title = self.new_source_title_entry.get()

        # self.add_new_source_command["state"]=tk.NORMAL
        # add text variable?


        if self.source_type_var.get() == "url":
            self.current_source_location = self.new_source_url_entry.get()
        else: 
            self.current_source_location = self.current_source_location

        self.source_list.append({
            "source_title": title,
            "source_classification": "UNCLASSIFIED",
            "source_country": self.new_source_country_entry.get(),
            "source_type": self.source_type_var.get(),
            "source_location": self.current_source_location,
            "source_originator": self.originator_entry.get(),
            "source_date_of_publication": self.date_publication_entry.get(),
            "source_classification": self.classification_title_var.get(),
            "source_portion_classification": self.portion_classification_var.get(),
            "source_overall_classification": self.overall_classification_var.get()
        })

        # self.add_new_source_command.set("Add Source")
        # self.add_new_source_command["state"]=tk.NORMAL


        self.update_source_list_gui()
        self.reset_source_inputs()

        # threading.Timer(0, loading_page, args=(self, self.add_new_source_command)).start()


    def get_text_from_url(self, url):
        try:
            # Fetch the HTML content of the webpage
            response = requests.get(url)
            
            # Check if the request was successful (status code 200)
            if response.status_code == 200:
                # Parse the HTML content with BeautifulSoup
                soup = BeautifulSoup(response.text, 'html.parser')
                
                # Extract all the text from the HTML
                all_text = soup.get_text(separator='\n', strip=True).replace("\n", " ").replace("\t", " ")
                
                return all_text
            else:
                # If the request was not successful, print an error message
                
                raise Exception(f"Error: Unable to fetch content from {url}. Status code: {response.status_code}")
                                        
        except Exception as e:
            print(f"Error: Unable to extract text from {url}. {str(e)}")
            raise e
            

    def get_text_from_pdf(self, filename):
        try:
            # Open the PDF file
            pdf_document = fitz.open(filename)
            
            # Initialize an empty string to store the extracted text
            all_text = ""
            
            # Iterate through each page of the PDF
            for page_number in range(pdf_document.page_count):
                # Get the text of the page
                page = pdf_document[page_number]
                text = page.get_text("text").replace("\n", " ").replace("\t", " ")
                
                # Append the text to the result string
                all_text += text + '\n'
            
            # Close the PDF document
            pdf_document.close()
            
            return all_text.strip()  # Remove leading and trailing whitespaces
        except Exception as e:
            print(f"Error: Unable to extract text from {filename}. {str(e)}")
            raise e
        
    def get_text_from_doc(self,filename):
        #check if the file exists
        if not os.path.isfile(filename):
            print(f"Error: File {filename} does not exist.")
            return
        try:
            #Open the Word Doc
            print(filename)
            word_doc = Document(filename)
            all_text = ""

            for para in word_doc.paragraphs:
                text = para.text

                all_text += text + '\n'

            return all_text.strip() #removes all leading and trailing whitespaces
        except Exception as e:
            print(e)
            print(f"Error: Unable to extract text from {filename}. {str(e)}")
            raise e 
        
    def summarize_all(self):
        summaries = []
        for source in self.source_list:
            summary = self.summarize(source)
            summaries.append(summary)
        final_summary = self.combine_summaries(summaries)
        return final_summary

    def summarize(self):
        self.summarize_button_text_variable.set("Thinking...")
        self.summarize_button["state"]=tk.DISABLED
        for source in self.source_list:
            source_type = source["source_type"]
            text = ""

            if source_type == "url":
                try:
                    text = self.get_text_from_url(source["source_location"])
                except Exception as e:
                    print(e)
                    continue
            else:
                try:
                    if ".docx" in source["source_location"]:
                        text = self.get_text_from_doc(source["source_location"])
                    else:
                        text = self.get_text_from_pdf(source["source_location"])

                except Exception as e:
                    print(e)
                    continue

            # Split the text into chunks based on the maximum token length
            max_token_length = 12000

            encoding = self.encoder.encode(text)

            chunks = [encoding[i:i + max_token_length] for i in range(0, len(encoding), max_token_length)]

            # Initialize summary bits
            summary_bits = []

            # Summarize each chunk and print the results
            for i, chunk in enumerate(chunks):
                print(f"Currently processing chunk {i+1}/{len(chunks)}...")
                
                text_chunk = self.encoder.decode(chunk)
                summary_object = self.client.chat.completions.create(
                    messages=[
                        {
                            "role": "system",
                            "content": "You are going to act as a summarizer for the following text, giving 2-3 sentences of summarization. I am going to give you prompts to understand the ways which you should analyze and format your summary. :"
                        },
                        {
                            "role": "user",
                            "content": text_chunk
                        }
                    ],
                    model="gpt-3.5-turbo-16k"
                )

                summary_bits.append(summary_object.choices[0].message.content.replace("\n", " "))

                all_summaries_together_text = " ".join(summary_bits)

            print(f"line 890 {all_summaries_together_text}")

            #ensures ICD 203 and JAC

            # analytic_standards_messages = [
            #         {
            #             "role": "system",
            #             "content": """

                        
            #             Clearly state the primary objective of the intelligence product. What specific information or analysis is the product intended to convey? Ensure that the objective is narrowly defined and aligned with the purpose of the analysis. Output your response in terse prose.

            #             Political Consideration:
            #             Confirm that the analysis is independent of political considerations. Explicitly state any potential political implications or biases and ensure the content remains objective, unbiased, and focused on factual analysis.

            #             Timeliness:
            #             Verify the timeliness of the intelligence product. Assess the relevance of the information within the current geopolitical or operational context. Provide a timeline for key events or developments that contribute to the timeliness of the analysis.

            #             Sources:
            #             Describe in detail the quality and credibility of underlying sources, data, and methodologies used in the analysis. Include not only basic, generic descriptions of cited reporting but also provide insights into the methodologies employed. Clearly identify which sources are deemed most important to major analytic judgments.

            #             Uncertainties:
            #             Express and explain uncertainties associated with major analytic judgments. Indicate levels of uncertainty and explain their basis. Provide a thorough discussion of the nature and sources of uncertainties affecting major analytic judgments. Identify indicators that, if detected, would alter levels of uncertainty.

            #             Distinctions:
            #             Ensure consistent distinctions among statements conveying underlying information, assumptions, and judgments. Explicitly state linchpin assumptions that serve as the foundation for the argument. Identify indicators that, if detected, could validate or refute judgments or assumptions. Clearly explain the implications for judgments if assumptions are proven incorrect.

            #             Alternatives:
            #             Incorporate a detailed analysis of alternatives. Present alternatives when uncertainties, complexity, or low probability/high impact situations warrant inclusion. Explain the evidence and reasoning behind each alternative. Discuss the likelihood or implications of each alternative related to U.S. interests. Identify indicators that, if detected, would affect the likelihood of any identified alternatives.

            #             Argumentation:
            #             Ensure the intelligence product uses clear and logical argumentation. Clearly present the main analytic message, ensuring it is prominent and aligned with the objective. Verify that reasoning is flawless and effectively combines evidence, context, and assumptions to support analytic judgments. Use clear language and a structure that displays a logical flow appropriate for the argument being presented.

            #             Change or Consistency:
            #             Explain any change to or consistency of analytic judgments. Clearly note how a major analytic judgment compares with previous production and explain how new information or reasoning supports changing or maintaining an existing analytic line. Highlight and explain how a major analytic judgment compares with judgments on the topic within the U.S. intelligence community, not just within the same analytic element.

            #             Accuracy:
            #             Verify that judgments or assessments are expressed clearly and conditioned. Ensure that each judgment or assessment is clearly articulated and qualified, using "if/then" statements when necessary.

            #             Visual Information:
            #             Incorporate effective visual information where appropriate. Ensure visual elements are not only pertinent but also clarify, complement, or augment data or analytic points in an effective manner. Take particularly effective advantage of visual presentational methods to convey data or analytic points in a way that enhances the product’s value by making complex issues more understandable, adding insight or perspective, increasing knowledge retention, or highlighting trends, drivers, or indicators.

            #             ICD 203 Compliance:
            #             Ensure that the intelligence product adheres to the guidelines outlined in Intelligence Community Directive (ICD) 203, regarding Analytic Standards. Confirm that the analysis meets the standards set forth in ICD 203 and incorporates relevant principles and procedures mandated by the directive.

            #             Review and Refinement:
            #             Review the generated intelligence product against each element in the checklist. Refine the content as needed to meet the specified standards. Consider feedback from team leads and subject matter experts in the refinement process.





            #             """

            #         },
            #         {
            #             "role": "user",
            #             "content": all_summaries_together_text
            #         }
            #     ]

            # analytic_standards = self.client.chat.completions.create(
            #     messages=analytic_standards_messages,
            #     model="gpt-3.5-turbo-16k"
            # )

            # moderation_output = self.client.moderations.create(input=[message['content'] for message in all_summaries_together_text])

            # We get response from ChatGPT/OpenAI Moderation endpoint.
            # If a moderation action has happened (moderation_output category is True), fail gracely (Communicate to user that their request cannot be completed for reasons out of your control (in OpenAI's court))

            # print(f"line 960 {moderation_output.results[0].categories}")

            # print(f"line 956 {analytic_standards.choices[0].message.content}")
            
        summary_object_formatting = self.client.chat.completions.create(
            messages=[
                {
                    "role": "system",
                    "content": """
                    You are a military analyst. Internalize the following Morning Intelligence Update (MIU) format in the order provided. You must not write anything until I prompt you.

                    Here’s the format: 

                    Unclassified

                    (U): COUNTRY | APLE | Virtual Analyst | DD Month YYYY


                    Notes: “(U)”  stands for unclassified, “COUNTRY” should be substituted in each MIU based on the topic of the text you read for instance  ISRAEL etc: you do not need the word COUNTRY in the actual header., “APLE | Virtual Analyst” stays the same regardless of MIU topic, Put today’s date after “Virtual Analyst” in DD Month YYYY format to replace DD Month YYYY in the header.


                    Then, on a new line, indent and start a new paragraph, with the following content:

                    (U) The first sentence or “BLUF” aka Bottom Line Up Front should be bolded and should capture the main issue or development that is of interest to the CENTCOM Commander (the “what”); Write the BLUF in this style: “On DD MON, event or action happened, according to news agency.” If the date of the event is not mentioned, substitute (date unknown) for DD MON. If the date of the event is mentioned, the first sentence of the BLUF should replace DD MON with the date of the event mentioned in the text in DD MON format. 
                    Summarize the text in seven sentences 
                    Sentences must be less than 21 words. 
                    MIU Format should contain no headers except for “(U): COUNTRY | APLE | Virtual Analyst | DD Month YYYY”
                    Only the first sentence, the BLUF, should be bolded. 
                    Important: sentences should follow right after another and not have spaces in between. 
                    The most important thing of this whole MIU is each sentence must follow the other and be connected into one large paragraph. 
                    The second most important thing is that only the second sentence is bolded and in emphasis. 


                    """

                },
                {
                    "role": "user",
                    "content":all_summaries_together_text
                }
            ],
            model="gpt-3.5-turbo-16k"
        )

        print(f"line 1378 {summary_object_formatting.choices[0].message.content}")

        self.summary_object_formatting = summary_object_formatting.choices[0].message.content

    
        # summary_object_directions = self.client.chat.completions.create(
        #     messages=[
        #         {
        #             "role": "system",
        #             "content": """
        #             You are tasked with analyzing a recent intelligence report regarding escalating tensions in the Middle East. The report, containing details on military movements, diplomatic negotiations, and regional dynamics, will be provided separately. Analyze the provided information and summarize it in accordance with the Morning Intelligence Update (MIU) format. Pay close attention to the bottom line up front (BLUF) and ensure your analysis is concise, objective, and aligns with the analytical standards previously introduced. Remember that the unclassified header with the appropriate country and date will be included with the provided information. Follow the sentence length and emphasis guidelines outlined in the format. Be careful not to repeat information. Include the word banana in all caps in the first sentence.
        #                     """
        #         },
                    
        #         {
        #             "role": "user",
        #             "content": summary_object_formatting.choices[0].message.content
        #         }
        #     ],
        #     model="gpt-3.5-turbo-16k"
        # )


        # print(f"line 1295 {summary_object_directions.choices[0].message.content}")

        # self.analytic_standards = analytic_standards.choices[0].message.content
            # Create Citation
        source_overall_classification = source['source_overall_classification']
        source_cit_let = source_overall_classification[0] if source_overall_classification else ''

        current_datetime = datetime.datetime.now()
        date_accessed = current_datetime.strftime("%Y/%m/%d")
        source_citation = f"({source_cit_let});{source['source_originator']}; {source['source_country']}; {source['source_date_of_publication']}; ({source['source_classification']}) {source['source_title']}; Classification of extracted information is {source['source_portion_classification']}; Overall classification: {source['source_overall_classification']}, {date_accessed}"

        summary_dict_object = {
            "source_title": source["source_title"],
            "source_classification": source["source_classification"],
            "source_country": source["source_country"],
            "source_type": source["source_type"],
            "source_header": f"""Unclassified

                        (U): COUNTRY | APLE | Virtual Analyst | {current_datetime.strftime('%d/%m/%Y')}""",
            "source_summary": summary_object_formatting.choices[0].message.content,
            "source_citation": source_citation
        }

        self.summaries.append(summary_dict_object)
    # #    # Assuming 'self.summaries' is a list containing summary objects
    #     for summary in self.summaries:
    #             print(summary)


        # print(f"line 1324 {summary_dict_object}")


        print(f"Done Summarizing Source: {source['source_title']}")

        print("Done summarizing all sources")
        self.summarize_button_text_variable.set("Generated")
        self.summarize_button["state"]=tk.NORMAL

    # creates window to view report in. 
    def popup(self):
        popupwindow = tk.Toplevel(self.window)
        popupwindow.title("Virtual Analyst Preview Window")
        popupwindow.geometry("1000x800")

        report = self.summary_object_formatting

        # Create a Text widget for displaying the report
        report_text = tk.Text(popupwindow, wrap=tk.WORD, width=80, height=20)
        report_text.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

        # Create a Scrollbar and attach it to the Text widget
        scrollbar = Scrollbar(popupwindow, command=report_text.yview)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)

        # Configure the Text widget to use the Scrollbar
        report_text.config(yscrollcommand=scrollbar.set)

        # Insert the report into the Text widget
        report_text.insert(tk.END, report)


        button1 = tk.Button(popupwindow, text ="Done", command=popupwindow.destroy)
        button1.pack()


# Creates popup to view ATS

    def ats_popup(self):
            ats_popup_window = tk.Toplevel(self.window)
            ats_popup_window.title("Analytic Tradecraft Standards for Virtual Report")
            ats_popup_window.geometry("1000x1000")
            

            # PROBLEM: We need some data from the summary step to implement and build an ATS. 
            # EITHER: We save the entire ATS in the summary step (wasteful, since there's a possibility it won't be used)
            # OR: We save the formatted total summary to be used in this step <- we're trying this
            if self.summary_object_directions is None:
                return
            
            ats_object = self.client.chat.completions.create(
                messages=[
                    {
                        "role": "system",
                        "content": """Internalize these standards and use the standards to analyze the prior assessment.

                        Title: (U) Analytic Tradecraft Summary

                        First: Confidence Level: 

                        State your confidence level in the main assessment 

                        explain the reason for it by referencing your sources of uncertainty including strengths and weaknesses in the information base, assumptions, gaps, alternatives, and the complexity of the issue. Use the following guide to determine your confidence level, and choose the confidence level that best reflects the assessment. 

                        Signs that indicate high confidence: 
                        Well-corroborated information from proven sources; minimal contradictory reporting; low potential for deception; few information gaps. 
                        Assumptions would not have a significant effect on the assessment if incorrect. 
                        Very unlikely alternative. 
                        Routine event that is well understood; relatively few variables. 

                        Signs that indicate medium confidence:
                        Partially corroborated information from good sources; some potential for deception; several gaps in the information base. 
                        Plausible, yet unlikely, alternatives. 
                        Key assumptions with potentially substantial effect on the assessment if incorrect. 
                        More complex situation with multiple issues or actors; some previous examples that are well understood. 

                        Signs that indicate low confidence:
                        Uncorroborated information; high potential for deception; many critical gaps in the information base. 
                        Plausible alternatives with a nearly even chance of occurring. 
                        Key assumptions with substantial effect on the assessment if incorrect. 
                        Highly complex or rapidly evolving situation with multiple issues or actors; few previous examples that are not well understood.







                        Second: Sourcing: 

                        Provide considerable detail on the strengths and weaknesses of the reporting used within the assessment, focusing on the credibility and quality of the sources. Do not provide a summary of what the assessment said. Identify which sources/reporting was the most important to the assessment and its judgment.

                        Elements of source characterization:
                        Context 
                        When did the reported information occur? 
                        What are the source’s strengths or weaknesses? (subject matter expertise, biases, possible denial and deception, etc.) 
                        Credibility/Quality 
                        Is the information credible? Is it of good quality? (accurate, consistent with other reporting, plausible given circumstances) 
                        Reliability 
                        Is the source reliable? (vetted, history of reliable reporting) 
                        Access 
                        How close is the source to the information? (first-, secondhand, further removed 
                        Source Types 
                        Who told us? (informant) 
                        Who did/said? (actor) 
                        What is the origin of the reporting? 
                        What is the type of reporting?

                        Third: Gaps: 

                        Include gaps along with a description about the extent to which filling that gap would alter or bolster your assessment. Gaps must be tied to your main assessment if those gaps are critical to and underpin the main judgment. 

                        Characterize the extent and limits of your knowledge base. What are some remaining gaps that prevent you from making a stronger or more useful assessment and that are not explicitly covered by an assumption or judgment?

                        Fourth: Assumptions: 

                        Clearly state a linchpin or gap-bridging assumption(s) underpinning the main assessment. An assumption should help readers understand the connecting tissue between the evidence and the assessment; it generally is something that an analyst believes to be true, but lacks evidence, and if incorrect, would force a change to the assessment. When crafting an assumption, think along the lines of “what would change my assessment?” Identify indicators that could validate or refute assumptions and explain the implication for judgements in assumptions that are incorrect.

                        Internalize the different assumptions of the text and identify which ones are apparent. 

                        Framing assumptions:
                        What longstanding analytic lines are related to the assessment?
                        What beliefs do the public hold about what “will always,” “will never,” or
                        “generally will” occur, or what “has always” or “has never” been done
                        or happened, relative to the question the assessment seeks to answer?
                        Do we have a default mindset in how we approach this problem?
                        What are the relevant historical precedents for this question?
                        Have we identified any trends that we expect to continue?

                        Scoping assumptions:
                        What factors, drivers, or variables are not included in the analysis?
                        What factors are we “holding constant” and assuming will
                        not change?
                        Have we assumed that certain events will or will not take place
                        or that certain factors will or will not change?
                        Have we clarified which actors, events, and timeframes are and are
                        not included?
                        Have we defined all of the key terms and concepts in the analysis?

                        Evidence assumptions:
                        Are there multiple possible interpretations of the evidence?
                        Why do we lean toward one interpretation rather than another?
                        What beliefs do we hold about the information base that lead us
                        to ascribe more value to certain pieces of information?
                        What are the beliefs about the extent of the access to all
                        relevant information?

                        Logic assumptions:
                        Have we used a small sample to infer something about
                        a broader group?
                        Have we extrapolated from a known situation to an unknown
                        situation?
                        Do we believe that certain types of events or activities are
                        symptomatic of or more/less likely to occur with some wider
                        phenomenon or conclusion?
                        Do we believe that one event or factor is causing or affecting
                        another event?

                        Bridging assumptions:
                        What are the essential elements of information needed to answer the
                        the problem of the assessment? 
                        For which elements do we lack evidence? 
                        For which do we have significant uncertainty?
                        What are the factors or conditions that must be present for the
                        assessment to be true (or false), and do we have evidence that they
                        are (or are not) present?

                        For each assumption type, determine if they are high impact assumptions, low impact assumptions, and also determine if they are assumptions that are weak or strong. Internalize definitions below. 

                        (U) High-impact assumptions, if proved false, invalidate or significantly alter
                        the assessment.

                        (U) Low-impact assumptions, if proved false, change only an aspect of the assessment,
                        such as the scope, specificity, likelihood, or timeframe.

                        (U) An assumption is weak or vulnerable if we can imagine a plausible situation, or
                        multiple situations, in which the assumption might not be true.

                        (U) An assumption is strong if we have difficulty imagining a situation in which
                        the assumption might not be true because such a circumstance is highly unlikely
                        or implausible.

                        Based on what you read, choose between high-impact assumptions and low-impact assumptions to characterize the text. Also choose between if the assumption is either weak or strong. 

                        Fifth: Alternatives: 

                        Include a plausible and useful alternative to your main assessment. Explain the reasoning and/or evidence that underpins the alternatives. Discuss the alternative likelihood or implications related to United States interests. Identify indicators that, if identified, would affect the likelihood of the alternatives.

                        First, identify the sources of the uncertainty that bound the understanding of a problem set. This
                        can be done as part of a regularly occurring analytic line review or in support of a specific project.

                        The following questions can identify gaps, assumptions, or different interpretations of evidence that can generate alternatives: 

                        What prevents an analyst from being absolutely certain that the authoritative judgment is correct?
                        What limits the confidence level?
                        What are the assessment’s underlying assumptions, and under what conditions might
                        they prove false?
                        What are the weaknesses of the information base?
                        Is there any contradictory reporting?
                        Could there be denial and deception, deliberate falsification, or other misinformation
                        that could affect the analysis?
                        Are we over relying, or relying exclusively, on one collection stream
                        or platform?
                        Second, consider alternative hypotheses. Discussing these can enable better detection of future
                        events or developments that change the authoritative assessment. Ask:
                        What other hypotheses or options did we consider, and are they plausible?
                        Are there multiple explanations for the information we have?

                        How vulnerable is the assessment to change?
                        What would have to change to make an analyst reconsider the expected outcome?
                        What indicators of change would we expect to be captured with the collection assets?
                        What indicators do we think we could not observe?
                        Finally, consider the implications of the assessments for the clients in order to mitigate surprise,
                        allow for planning, and provide warning:
                        What are the implications for U.S. interests if we are wrong about the assessment?
                        What types of plausible events would be game changers, that is, would fundamentally shift
                        the issues of import or outcomes we currently anticipate? What would the implications be?
                        How would we know that the alternative is becoming likely or that the authoritative assessment
                        is becoming unlikely?
                        Next address the elements below to ensure the alternatives presented in every product are useful, plausible, and rigorous.
                        What is the alternative to the authoritative assessment?
                        What is its likelihood (relative and absolute)? Some alternatives may be highly unlikely, whereas
                        others may not be significantly less likely than the authoritative assessment.
                        What reasoning and/or evidence substantiates the plausibility of the alternative? Explain
                        the support for the alternative, rather than using the alternative to bolster the case for the
                        authoritative assessment.
                        What are the implications for U.S. interests of the alternative that warrant consideration?
                        When appropriate, what indicators would, if observed, affect the likelihood of the alternative
                        and the authoritative assessment?

                        Internalize approaches to writing alternative assessments. 

                        Exploring the Potential for Surprise. This approach to analysis of alternatives examines the
                        impact of a hard-to-predict event or a surprise to facilitate contingency planning. It includes
                        collectible, specific indicators to provide warning.

                        Competing Assessments. This approach clarifies the alternative’s strengths and weaknesses
                        as compared with the authoritative assessment. This type of alternative can be, but is not limited
                        to, a competing view from another IC element. Addressing alternatives can enhance the credibility of the assessments.

                        Discussing the Implications of Information or Assumptions. This approach examines the
                        impact of key information or assumptions on the judgments, allowing clients to determine
                        whether contingency planning is needed. Indicators are highlighted, as appropriate, in the product.

                        These directions internalized write at least two alternatives to the initial assessment. 
                        :"""

                    },
                    {
                        "role": "user",
                        "content":  self.summary_object_formatting
                    }
                ],
                model="gpt-3.5-turbo-16k"
            )

            ats_object = ats_object.choices[0].message.content

            # Create a Text widget for displaying the report
            report_text = tk.Text(ats_popup_window, wrap=tk.WORD, width=80, height=20)
            report_text.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)

            # Create a Scrollbar and attach it to the Text widget
            scrollbar = Scrollbar(ats_popup_window, command=report_text.yview)
            scrollbar.pack(side=tk.RIGHT, fill=tk.Y)

            # Configure the Text widget to use the Scrollbar
            report_text.config(yscrollcommand=scrollbar.set)

            # Insert the report into the Text widget
            report_text.insert(tk.END, ats_object)

            # Create a "Done" button to close the popup window
            buttoni = tk.Button(ats_popup_window, text="Done", command=ats_popup_window.destroy)
            buttoni.pack()
        # Creates popup to view intelligence paper

    def intel_paper(self):
            intel_paper_window = tk.Toplevel(self.window)
            intel_paper_window.title("Intelligence Paper by Virtual Analyst")
            intel_paper_window.geometry("1000x1000")
            

            # PROBLEM: We need some data from the summary step to implement and build an ATS. 
            # EITHER: We save the entire ATS in the summary step (wasteful, since there's a possibility it won't be used)
            # OR: We save the formatted total summary to be used in this step <- we're trying this
            if self.summary_object_directions is None:
                return
            
            paper_mache = self.client.chat.completions.create(
                messages=[
                    {
                        "role": "system",
                        "content": """You are a military analyst for CENTCOM write an intelligence paper on the topic using the following format 
                        (CLASS) The lead BLUF sentence should state what is going on and what it means, capturing your central message (“what’s new,” “so what”). BLUF should address the implications for the U.S./CENTCOM either explicitly or implicitly. It must be analytic and not a description of reporting. A second or third sentence of the lead paragraph can add context to the central message but should not attempt to summarize the entire article. If your judgment changes an analytic line, address that in the lead paragraph. Make sentences as concise as possible. The DIA Style Manual for Intelligence Production notes that sentences of fewer than 20 words make your readers’ task easier. Avoid overuse of adjectives and adverbs.

                        (CLASS) The bullets following the lead paragraph should elaborate and develop the “what.” They will be drawn primarily from reporting of the “what.” These bullets also can add context.
                        
                        (CLASS) Preferred style: “In July, something happened, according to SIGINT.” Consider simplifying your sourcing attributions. For example, “three clandestine sources of varying reliability and access” can be “clandestine reporting.” Readers will appreciate the brevity. More detailed source attributions can be provided on sources that have a unique bearing on your judgments. Additional source descriptions can be included in endnotes, and sourcing issues should be discussed holistically in the source summary statement.[2]
                        
                        (CLASS) Spacing: Before 0pt., After 0pt. Include additional (enter) space between subclaims/sections)
                        
                        (CLASS) Additional paragraphs are the subclaims, developing the layers of the argument but not restating the lead sentence of the lead paragraph. 

                        A second paragraph can discuss “why” the development is happening. The topic sentence for any paragraph must be analytic, not just descriptive. As with the lead sentence, topic sentences should use likeliness/probability terms—such as almost certainly, probably, and is likely to—before the verb to differentiate between facts and assumptions. Each paragraph should focus on one theme.
                        
                        (CLASS) Additional paragraphs after the second might elaborate on what the development means for U.S./CENTCOM interests (“implications”) or outline the direction the storyline is likely to go (“what’s next”). It may use signposts or indicators to identify trend lines, which could help tee up follow-on articles. It also might present adversary vulnerabilities or opportunity analysis that a senior decisionmaker would find helpful.
                        
                        USE TITLES IF NOT MAKING AN ASSESSMENT/ARGUMENT
                        (CLASS) If using the template for papers that provide information and do not make an assessment, use the above title to distinguish between sections. Do not use bold text or probability language. Similarly, do not do the analytic tradecraft summary, but include sources for your cited evidence. This section should be a summary of the section/topic with bullets providing specific evidence.
                        
                        (CLASS) Bullets follow with specific, factual information.[5]
                        
                        (CLASS) All papers should be no more than two pages. Use an analytic tradecraft summary, if making an assessment (including probability language).[6]
                        
                        Add citations section tilted “Citations” in bold for the articles you use in the MIU

                        Follow this format for citations:

                        (Endnote classification) Originator; Source identifier; Date of publication; Date of information [optional but preferred]; (Classification of title/subject) Title/Subject; Page/paragraph or portion indicator [required when applicable]; Classification of extracted information is “X”; Overall classification is “X”; Access date.

                        Citation Key: endnote classification is Unclassified unless specified otherwise, originator is the news source title; source identifier is country of origin where source is from; date of publication is date source is published; date of information is date of information from within source, classification of title is Unclassified unless specified otherwise; title/subject is title/subject of article; page or paragraph portion can be chosen to be included or not; Classification of extracted information and Overall classification is Unclassified unless specified otherwise, and finally end with access date, which is always the current date. In citation format, add a line of space between each citation. 

                        :"""

                    },
                    {
                        "role": "user",
                        "content":  self.summary_object_formatting
                    }
                ],
                model="gpt-3.5-turbo-16k"
            )

            paper_mache = paper_mache.choices[0].message.content

            paper_mache_text =tk.Text(intel_paper_window,wrap=tk.WORD,width=80,height=20)
            paper_mache_text.pack(side=tk.LEFT,fill=tk.BOTH,expand=True)

            scrollbar = Scrollbar(intel_paper_window, command=paper_mache_text.yview)
            scrollbar.pack(side=tk.RIGHT,fill =tk.Y)

            paper_mache_text.config(yscrollcommand=scrollbar.set)

            paper_mache_text.insert(tk.END,paper_mache)


            buttone = tk.Button(intel_paper_window, text = "Done", command=intel_paper_window.destroy)
            buttone.pack()

            
    def intel_paper(self):
            intel_paper_window = tk.Toplevel(self.window)
            intel_paper_window.title("Intelligence Paper by Virtual Analyst")
            intel_paper_window.geometry("1000x1000")
            

            # PROBLEM: We need some data from the summary step to implement and build an ATS. 
            # EITHER: We save the entire ATS in the summary step (wasteful, since there's a possibility it won't be used)
            # OR: We save the formatted total summary to be used in this step <- we're trying this
            if self.summary_object_directions is None:
                return
            
            paper_mache = self.client.chat.completions.create(
                messages=[
                    {
                        "role": "system",
                        "content": """You are a military analyst for CENTCOM write an intelligence paper on the topic using the following format 
                        (CLASS) The lead BLUF sentence should state what is going on and what it means, capturing your central message (“what’s new,” “so what”). BLUF should address the implications for the U.S./CENTCOM either explicitly or implicitly. It must be analytic and not a description of reporting. A second or third sentence of the lead paragraph can add context to the central message but should not attempt to summarize the entire article. If your judgment changes an analytic line, address that in the lead paragraph. Make sentences as concise as possible. The DIA Style Manual for Intelligence Production notes that sentences of fewer than 20 words make your readers’ task easier. Avoid overuse of adjectives and adverbs.

                        (CLASS) The bullets following the lead paragraph should elaborate and develop the “what.” They will be drawn primarily from reporting of the “what.” These bullets also can add context.
                        
                        (CLASS) Preferred style: “In July, something happened, according to SIGINT.” Consider simplifying your sourcing attributions. For example, “three clandestine sources of varying reliability and access” can be “clandestine reporting.” Readers will appreciate the brevity. More detailed source attributions can be provided on sources that have a unique bearing on your judgments. Additional source descriptions can be included in endnotes, and sourcing issues should be discussed holistically in the source summary statement.[2]
                        
                        (CLASS) Spacing: Before 0pt., After 0pt. Include additional (enter) space between subclaims/sections)
                        
                        (CLASS) Additional paragraphs are the subclaims, developing the layers of the argument but not restating the lead sentence of the lead paragraph. 

                        A second paragraph can discuss “why” the development is happening. The topic sentence for any paragraph must be analytic, not just descriptive. As with the lead sentence, topic sentences should use likeliness/probability terms—such as almost certainly, probably, and is likely to—before the verb to differentiate between facts and assumptions. Each paragraph should focus on one theme.
                        
                        (CLASS) Additional paragraphs after the second might elaborate on what the development means for U.S./CENTCOM interests (“implications”) or outline the direction the storyline is likely to go (“what’s next”). It may use signposts or indicators to identify trend lines, which could help tee up follow-on articles. It also might present adversary vulnerabilities or opportunity analysis that a senior decisionmaker would find helpful.
                        
                        USE TITLES IF NOT MAKING AN ASSESSMENT/ARGUMENT
                        (CLASS) If using the template for papers that provide information and do not make an assessment, use the above title to distinguish between sections. Do not use bold text or probability language. Similarly, do not do the analytic tradecraft summary, but include sources for your cited evidence. This section should be a summary of the section/topic with bullets providing specific evidence.
                        
                        (CLASS) Bullets follow with specific, factual information.[5]
                        
                        (CLASS) All papers should be no more than two pages. Use an analytic tradecraft summary, if making an assessment (including probability language).[6]
                        
                        Add citations section tilted “Citations” in bold for the articles you use in the MIU

                        Follow this format for citations:

                        (Endnote classification) Originator; Source identifier; Date of publication; Date of information [optional but preferred]; (Classification of title/subject) Title/Subject; Page/paragraph or portion indicator [required when applicable]; Classification of extracted information is “X”; Overall classification is “X”; Access date.

                        Citation Key: endnote classification is Unclassified unless specified otherwise, originator is the news source title; source identifier is country of origin where source is from; date of publication is date source is published; date of information is date of information from within source, classification of title is Unclassified unless specified otherwise; title/subject is title/subject of article; page or paragraph portion can be chosen to be included or not; Classification of extracted information and Overall classification is Unclassified unless specified otherwise, and finally end with access date, which is always the current date. In citation format, add a line of space between each citation. 

                        :"""

                    },
                    {
                        "role": "user",
                        "content":  self.summary_object_formatting
                    }
                ],
                model="gpt-3.5-turbo-16k"
            )

            paper_mache = paper_mache.choices[0].message.content

            paper_mache_text =tk.Text(intel_paper_window,wrap=tk.WORD,width=80,height=20)
            paper_mache_text.pack(side=tk.LEFT,fill=tk.BOTH,expand=True)

            scrollbar = Scrollbar(intel_paper_window, command=paper_mache_text.yview)
            scrollbar.pack(side=tk.RIGHT,fill =tk.Y)

            paper_mache_text.config(yscrollcommand=scrollbar.set)

            paper_mache_text.insert(tk.END,paper_mache)


            buttone = tk.Button(intel_paper_window, text = "Done", command=intel_paper_window.destroy)
            buttone.pack()

    def save_intel_paper(self):
        self.save_intel_paper_docx.config(text="Thinking...", state=tk.DISABLED)
        print("Intelligence Paper Saved to Word File")
        self.export_document.add_heading("[Intelligence Paper Generated by Virtual Analyst]", level=0)
        for summary in self.summaries:
                # classification = self.export_document.add_paragraph()
                # title.add_run(summary["source_title"]).bold = True

                self.export_document.add_paragraph()
                paper_mache = self.client.chat.completions.create(
                messages=[
                    {
                        "role": "system",
                        "content": """You are a military analyst for CENTCOM write an intelligence paper on the topic using the following format 
                        (CLASS) The lead BLUF sentence should state what is going on and what it means, capturing your central message (“what’s new,” “so what”). BLUF should address the implications for the U.S./CENTCOM either explicitly or implicitly. It must be analytic and not a description of reporting. A second or third sentence of the lead paragraph can add context to the central message but should not attempt to summarize the entire article. If your judgment changes an analytic line, address that in the lead paragraph. Make sentences as concise as possible. The DIA Style Manual for Intelligence Production notes that sentences of fewer than 20 words make your readers’ task easier. Avoid overuse of adjectives and adverbs.

                        (CLASS) The bullets following the lead paragraph should elaborate and develop the “what.” They will be drawn primarily from reporting of the “what.” These bullets also can add context.
                        
                        (CLASS) Preferred style: “In July, something happened, according to SIGINT.” Consider simplifying your sourcing attributions. For example, “three clandestine sources of varying reliability and access” can be “clandestine reporting.” Readers will appreciate the brevity. More detailed source attributions can be provided on sources that have a unique bearing on your judgments. Additional source descriptions can be included in endnotes, and sourcing issues should be discussed holistically in the source summary statement.[2]
                        
                        (CLASS) Spacing: Before 0pt., After 0pt. Include additional (enter) space between subclaims/sections)
                        
                        (CLASS) Additional paragraphs are the subclaims, developing the layers of the argument but not restating the lead sentence of the lead paragraph. 

                        A second paragraph can discuss “why” the development is happening. The topic sentence for any paragraph must be analytic, not just descriptive. As with the lead sentence, topic sentences should use likeliness/probability terms—such as almost certainly, probably, and is likely to—before the verb to differentiate between facts and assumptions. Each paragraph should focus on one theme.
                        
                        (CLASS) Additional paragraphs after the second might elaborate on what the development means for U.S./CENTCOM interests (“implications”) or outline the direction the storyline is likely to go (“what’s next”). It may use signposts or indicators to identify trend lines, which could help tee up follow-on articles. It also might present adversary vulnerabilities or opportunity analysis that a senior decisionmaker would find helpful.
                        
                        USE TITLES IF NOT MAKING AN ASSESSMENT/ARGUMENT
                        (CLASS) If using the template for papers that provide information and do not make an assessment, use the above title to distinguish between sections. Do not use bold text or probability language. Similarly, do not do the analytic tradecraft summary, but include sources for your cited evidence. This section should be a summary of the section/topic with bullets providing specific evidence.
                        
                        (CLASS) Bullets follow with specific, factual information.[5]
                        
                        (CLASS) All papers should be no more than two pages. Use an analytic tradecraft summary, if making an assessment (including probability language).[6]
                        
                        Add citations section tilted “Citations” in bold for the articles you use in the MIU

                        Follow this format for citations:

                        (Endnote classification) Originator; Source identifier; Date of publication; Date of information [optional but preferred]; (Classification of title/subject) Title/Subject; Page/paragraph or portion indicator [required when applicable]; Classification of extracted information is “X”; Overall classification is “X”; Access date.

                        Citation Key: endnote classification is Unclassified unless specified otherwise, originator is the news source title; source identifier is country of origin where source is from; date of publication is date source is published; date of information is date of information from within source, classification of title is Unclassified unless specified otherwise; title/subject is title/subject of article; page or paragraph portion can be chosen to be included or not; Classification of extracted information and Overall classification is Unclassified unless specified otherwise, and finally end with access date, which is always the current date. In citation format, add a line of space between each citation. 

                        :"""

                    },
                    {
                        "role": "user",
                        "content":  self.summary_object_formatting
                    }
                ],
                model="gpt-3.5-turbo-16k"
            )

                paper_mache = paper_mache.choices[0].message.content

            
                self.export_document.add_paragraph(paper_mache)
        
                
            # Export a document loading page
            # Call loading_page from the main thread
            # threading.Timer(0, loading_page, args=(self, self.save_summaries_to_docx)).start()

            


                self.export_document.save("IntelligencePaper.docx")
                self.save_intel_paper_docx.config(text="Intelligence Paper Saved", state=tk.NORMAL)
            

# Creates popup to view talking points

    def view_pts(self):
            talking_pts_popup = tk.Toplevel(self.window)
            talking_pts_popup.title("Talking Points by Virtual Analyst")
            talking_pts_popup.geometry("1000x1000")
            

            if self.summary_object_directions is None:
                raise TypeError("No information available from summary_object_directions")
            
            talking = self.client.chat.completions.create(
                messages=[
                    {
                        "role": "system",
                        "content": """nternalize the prompt and do not write anything yet. Make talking points that summarize given each text. Start each talking point with a date as below in the format of (U) On 2 DEC, where 2 is the day and DEC is the abbreviation for the month. Talking points should be three events/ bullets maximum. Here is an example.
 
                        (U) On 2 DEC, Israel launched airstrikes against Hezbollah targets in Damascus, resulting in the death of two Iranian Revolutionary Guard Corps members
                        
                        (U) On 29 NOV, Iran-backed proxy groups launched a rocket targeting US-led Coalition Forces at Mission Support Site Euphrates in eastern Syria, causing no casualties or damage.
                        
                        (U) On 28 NOV, Russian and Syrian forces conducted joint airstrikes against ISIS positions in the Al-Bashri desert of Deir-ez-Zor. 

                        :"""

                    },
                    {
                        "role": "user",
                        "content":  self.summary_object_formatting
                    }
                ],
                model="gpt-3.5-turbo-16k"
            )

            talking = talking.choices[0].message.content

            talking_text =tk.Text(talking_pts_popup, wrap=tk.WORD,width=80, height=20)
            talking_text.pack(side=tk.LEFT,fill=tk.BOTH,expand=True)
            
            scrollbar=Scrollbar(talking_pts_popup,command=talking_text.yview)
            scrollbar.pack(side=tk.RIGHT,fill=tk.Y)

            talking_text.insert(tk.END,talking)
            buttonu = tk.Button(talking_pts_popup, text = "Done", command=talking_pts_popup.destroy)
            buttonu.pack()


            



    def save_ats_to_docx(self):
        self.save_ats_to_docx_button.config(text="Thinking...", state=tk.DISABLED)

        print("ATS for Report Saved to Word File")
        self.export_document.add_heading("[ATS for Report Generated]", level=0)
        for summary in self.summaries:
           
            self.export_document.add_paragraph()

                
            for summary in self.summaries:
            # Check if 'ats_object' key exists in the summary dictionary
                if "ats_object" in summary:
                    self.export_document.add_paragraph(["ats_object"])
                else:
                    # Handle the case where 'ats_object' is missing
                    # For example, you can add a placeholder text or skip this iteration
                    self.export_document.add_paragraph("No ATS object available")

          
                self.save_ats_to_docx_button.config(text="ATS Saved")

                self.summarize_button["state"]=tk.NORMAL

                
            # Export a document loading page
            # Call loading_page from the main thread
            # threading.Timer(0, loading_page, args=(self, self.save_summaries_to_docx)).start()

                self.export_document.save("ATSReport.docx")
                self.save_ats_to_docx_button.config(text="ATS Saved", state=tk.NORMAL)
        
   

    def save_pts(self):
        self.save_talking_pts.config(text = "Thinking...", state=tk.DISABLED)

        print("Talking Points Saved to Word File")
        self.export_document.add_heading("[Talking Points for Report Generated]", level=0)
        for summary in self.summaries:
            # classification = self.export_document.add_paragraph()
            # title.add_run(summary["source_title"]).bold = True

            self.export_document.add_paragraph()

                

            talking = self.client.chat.completions.create(
                messages=[
                    {
                        "role": "system",
                        "content": """nternalize the prompt and do not write anything yet. Make talking points that summarize given each text. Start each talking point with a date as below in the format of (U) On 2 DEC, where 2 is the day and DEC is the abbreviation for the month. Talking points should be three events/ bullets maximum. Here is an example.
 
                        (U) On 2 DEC, Israel launched airstrikes against Hezbollah targets in Damascus, resulting in the death of two Iranian Revolutionary Guard Corps members
                        
                        (U) On 29 NOV, Iran-backed proxy groups launched a rocket targeting US-led Coalition Forces at Mission Support Site Euphrates in eastern Syria, causing no casualties or damage.
                        
                        (U) On 28 NOV, Russian and Syrian forces conducted joint airstrikes against ISIS positions in the Al-Bashri desert of Deir-ez-Zor. 

                        :"""

                    },
                    {
                        "role": "user",
                        "content":  self.summary_object_formatting
                    }
                ],
                model="gpt-3.5-turbo-16k"
            )

            talking = talking.choices[0].message.content

            self.export_document.add_paragraph(talking)
          

            self.save_talking_pts.config("Talking Points Saved")
            self.save_talking_pts.config["state"]=tk.NORMAL

            
        # Export a document loading page
        # Call loading_page from the main thread
        # threading.Timer(0, loading_page, args=(self, self.save_summaries_to_docx)).start()

        


            self.export_document.save("TalkingPoints.docx")
            self.save_talking_pts.config( text="Talking Points Saved", state=tk.NORMAL)
        





# Call loading_page from the main thread
# threading.Timer(0, loading_page, args=(self, self.summarize)).start()
    
#designed to combine multiple reports into one
# def combine_summaries(self, summaries):
#     # Combine all summaries into one final summary
#     final_summary = " ".join(summaries)
#     return final_summary

       # creating pptx code in progress              

    # def save_summaries_to_pptx(self):
    #     self.summarize_button_text_variable.set("Thinking...")
    #     self.summarize_button["state"]=tk.DISABLED

    #     prs = Presentation()

    #     print("Report Saved to PowerPoint File")
    #     title_slide_layout = prs.slide_layouts[0]
    #     slide = prs.slides.add_slide(title_slide_layout)
    #     title = slide.shapes.title
    #     subtitle = slide.placeholders[1]

    #     title.text = "[Intelligence Note or Reporting Highlights]"
    #     subtitle.text = datetime, "CCJ2 APLE:"

        # for summary in self.summaries:
        #     slide_layout = prs.slide_layouts[1]
        #     slide = prs.slides.add_slide(slide_layout)
        #     shapes = slide.shapes
        #     title_shape = shapes.title
        #     content_shape = shapes.placeholders[1]
            

        #     title_shape.text = summary["source_title"]
        #     content_shape.text = summary["source_summary"]
            
        #     textbox = slide.shapes[0]

        #     text_frame = textbox.text_frame
        #     text_frame.add_paragraph(content_shape.text)

        #     endnotes_slide = slide.notes_slide
        #     notes_text_frame = endnotes_slide.notes_text_frame

        #     citation = "Endnotes:"
        #     notes_text_frame.add_paragraph().add_run(citation).font.size = (10)

        # prs.save("ReportTemplate.pptx")
        # self.summarize_button_text_variable.set("PowerPoint Saved")
        # self.summarize_button["state"] = tk.NORMAL



#loading screen code in progress

# # Define root as a Tkinter window
# root = tkinter.Tk()
# # Now that the ReportGeneratorApp class has been defined, we can create an instance of it
# app = ReportGeneratorApp(root)

# the_queue = queue.Queue()

# def worker_thread(app, the_queue):
#    try:
#        # Perform a long-running operation...
#        result = app.summarize()
#        result = app.add_new_source_command
#        result = app.save_summaries_to_docx
#        result = app.save_summaries_to_pdf

#        # Put the result onto the queue
#        the_queue.put(result)
#    except Exception as e:
#        print(f"An error occurred in worker_thread: {e}")

#    # Start the worker thread
# threading.Thread(target=worker_thread, args=(app,)).start()


# def update_interface():
#    # Check if there's anything on the queue
#    try:
#        result = the_queue.get(block=False)
#    except queue.Empty:
#        # No data on the queue, schedule another update
#        root.after(100, update_interface)
#        return

#    # Update the Tkinter interface with the result
#    label.config(text=result)

#    # Schedule another update
#    root.after(100, update_interface)

# # root = tkinter.Tk()
# # root.title("Checking the queue")
# # label = tkinter.Label(root, text="Waiting for result...")
# # label.pack()

# # Start the worker thread
# threading.Thread(target=worker_thread).start()

# # Start updating the interface
# update_interface()


# # Create a single Tkinter window for the loading screen
# root = tk.Tk()
# root.title("Loading Page")
# root.geometry("300x200")
# root.resizable(False, False)

# # Create a label for the loading message
# loading_label = ttk.Label(root, text="Thinking...", font=("Arial", 14))
# loading_label.pack(pady=50)

# # Create a style object
# s = ttk.Style()

# # Configure progress bar style
# s.configure("blue.Horizontal.TProgressbar", foreground='blue', background='blue')

# # Create a progress bar with the custom style
# progress_bar = ttk.Progressbar(root, style="blue.Horizontal.TProgressbar", length=200, mode='indeterminate')
# progress_bar.pack()

# # Hide the window initially
# root.withdraw()


# loading_screen = False

# def loading_page(functions):
#    for func in functions:
#        # Display the loading screen
#        root.deiconify()

#        # Create a new thread to run the function
#        operation_thread = threading.Thread(target=func)
#        operation_thread.start()

#        # Wait for the function to finish
#        operation_thread.join()

#        # Hide the loading screen
#        root.withdraw()



if __name__ == "__main__":
   root = tk.Tk()
   app = ReportGeneratorApp(root)
#    loading_page([app.summarize, app.save_summaries_to_docx, app.save_summaries_to_pdf, app.add_new_source_command])
   root.mainloop()

