from pprint import pp
from string import capwords
import tkinter as tk
from tkinter import ttk
from tkinter import *
from tkinter import Tk, Label, PhotoImage
from time import sleep
from tkinter import Canvas as tkCanvas, filedialog
from tkinter import messagebox
import tkinter
from xml.sax.handler import property_xml_string
from click import wrap_text
from transformers import pipeline, Pipeline, PreTrainedTokenizerFast, BartTokenizer,BartForConditionalGeneration
import requests
from bs4 import BeautifulSoup, Comment
import fitz
from openai import OpenAI
from dotenv import load_dotenv
import os
import tiktoken
from docx import Document
from reportlab.pdfgen.canvas import Canvas as pdfCanvas
from reportlab.lib.pagesizes import LETTER
from textwrap import wrap
import datetime
from threading import Thread
import pkg_resources
import threading
import queue
from pptx import Presentation
from tkinter import messagebox

# ats to pdf
# ATS
# ats prompt isnt being accessed -why? (line 823)
# link ats_object to use_ats_popup and save_ats_to_docx
# preview function
# find a rag
# joel - header?

# def is_installed(package_name):
#    try:
#        dist = pkg_resources.get_distribution(package_name)
#        print(f"{dist.key} ({dist.version}) is installed")
#        return True
#    except pkg_resources.DistributionNotFound:
#        print(f"{package_name} is NOT installed")
#        return False

# is_installed('python-docx')


#BG (blue = #6EB3F4), + image in python-first-steps, bg.  

tokenizer = BartTokenizer.from_pretrained('facebook/bart-large-cnn')
model = BartForConditionalGeneration.from_pretrained('facebook/bart-large-cnn')


load_dotenv()


class ReportGeneratorApp:
    def __init__(self, root):

        self.window = root
        # Initialize variables
        self.root: tk.Tk = root
        self.bg_image = PhotoImage(file="/Users/kendrafrench/Dev/python-first-steps/BG/verticaldblcrestltblbg.png")
        # create a label 
        bg_label = Label(root, image = self.bg_image)
        bg_label.place(x=0,y=0, relwidth=1,relheight=1)

        self.source_list: list[dict[str, str]] = []
        self.source_widgets: list[tk.Widget] = []
        self.title: str = "CCJ2-APLE Virtual Analyst"

        
        
       

        # Initialize source type variable
        self.source_type_var = tk.StringVar()
        

        # Initialize widgets
        self.root.title(self.title)
        self.create_initial_widgets()

        self.current_source_location: str = ""

        # # Initialize Model LEGACY, SWITCHED TO OPENAI
        # self.summarizer: Pipeline = pipeline("summarization", model="facebook/bart-large-cnn")
        
        # Initialize OpenAI
        self.client = OpenAI(
            api_key=os.getenv("SHIELDS_KEY"),
        )

        # Initialize encoder
        self.encoder = tiktoken.get_encoding("cl100k_base")

        self.export_document = Document()

        # Initialize Summaries
        self.summaries: list[dict[str, str]] = []

        # self.root = tk.Tk() 
        self.export_pdf = pdfCanvas("ReportTemplate-pdf", pagesize=LETTER)


    def save_summaries_to_pdf(self):
        # self.save_pdf_text_variable.set("Saving...")
        # self.save_pdf_text_variable["state"]=tk.DISABLED
        # need to add coding to the button itself

        print("Report saved to PDF")
        canvas = pdfCanvas("ReportTemplate.pdf", pagesize=LETTER)

        # heading formatting
        canvas.setFont("Helvetica", 18)
        canvas.drawString(72,740, "[Intelligence Note or Reporting Highlights]")
        canvas.setFont("Helvetica", 12)
            
        # formatting contents
        y = 700
        max_chars_per_line = 90 # Adjust this value as needed
        for summary in self.summaries:
            classification = summary["source_classification"]
            country = summary["source_country"]
            title = summary["source_title"]
            summary_text = summary["source_summary"]
            citation = summary["source_citation"]

            canvas.drawString(72, y, f"Classification: {classification}")
            y -= 20
            canvas.drawString(72, y, f"Country: {country}")
            y -= 20
            canvas.drawString(72, y, f"Title: {title}")
            y -= 40
            # canvas.drawString(72, y, f"{summary_text}")
            # split the text into lines
            summary_lines = wrap(summary_text, width=max_chars_per_line)
            for i, line in enumerate(summary_lines):
                canvas.drawString(72, y - i*20, line)
            y -= len(summary_lines)*20 + 60
            y -= 0

            canvas.drawString(72, y, f"Analyst Comment:")
            y = -20
            # split the text into lines
            citation_lines = wrap(citation, width=max_chars_per_line)
            for i, line in enumerate(citation_lines):
                canvas.drawString(72, y - i*20, line)
            y -= len(citation_lines)*20 + 60
            y -= 0


            # self.save_pdf_text_variable.set("PDF Saved")
            # self.save_pdf_text_variable["state"]=tk.NORMAL

            canvas.save()

            # Call loading_page from the main thread
            # threading.Timer(0, loading_page, args=(self, self.save_summaries_to_pdf)).start()

            

    def save_summaries_to_docx(self):
        self.summarize_button_text_variable.set("Thinking...")
        self.summarize_button["state"]=tk.DISABLED

        print("Report Saved to Word File")
        self.export_document.add_heading("[Intelligence Note or Reporting Highlights]", level=0)
        for summary in self.summaries:
            classification = self.export_document.add_paragraph()
            classification.add_run(summary["source_classification"]).bold = True
            # add caps to classification
            country = self.export_document.add_paragraph()
            country.add_run(summary["source_country"]).bold = True
            title = self.export_document.add_paragraph()
            title.add_run(summary["source_title"]).bold = True

            self.export_document.add_paragraph()

                
            self.export_document.add_paragraph(summary["source_summary"])
            
            add_analystc_comment = self.export_document.add_paragraph()

            self.export_document.add_paragraph()

            add_analystc_comment.add_run("[Analyst Comment]").bold = True

            citation = self.export_document.add_paragraph()
            citation.add_run(summary["source_citation"])


            self.summarize_button_text_variable.set("Document Saved")
            self.summarize_button["state"]=tk.NORMAL

            
        # Export a document loading page
        # Call loading_page from the main thread
        # threading.Timer(0, loading_page, args=(self, self.save_summaries_to_docx)).start()

        


            self.export_document.save("ReportTemplate.docx")
            self.summarize_button_text_variable.set("Summarize Sources")
            self.summarize_button["state"]=tk.NORMAL

    def save_ats_to_docx(self):
        self.save_ats_to_docx_button.set("Thinking...")
        self.save_ats_to_docx_button["state"]=tk.DISABLED

        print("ATS for Report Saved to Word File")
        self.export_document.add_heading("[ATS for Report Generated]", level=0)
        for summary in self.summaries:
            # classification = self.export_document.add_paragraph()
            # title.add_run(summary["source_title"]).bold = True

            self.export_document.add_paragraph()

                
            self.export_document.add_paragraph(summary["ats_object"])
          

            self.save_ats_to_docx_button.sete.set("ATS Saved")
            self.summarize_button["state"]=tk.NORMAL

            
        # Export a document loading page
        # Call loading_page from the main thread
        # threading.Timer(0, loading_page, args=(self, self.save_summaries_to_docx)).start()

        


            self.export_document.save("ATSReport.docx")
            self.save_ats_to_docx_button.set.set("Summarize Sources")
            self.save_ats_to_docx_button["state"]=tk.NORMAL




    def create_initial_widgets(self):
        # Create title
        title_label = tk.Label(self.root, text=self.title, font=("Arial", 25),bg="#6EB3F4")
        title_label.grid(column=0, row=0)

        # New Source Label
        self.source_list_label = tk.Label(self.root, text="Add Source", font=("Arial", 18), bg="#6EB3F4")
        self.source_list_label.grid(column=0, row=1)

        # Create first source frame
        new_source_frame = tk.Frame(self.root, bg="#6EB3F4")
        new_source_frame.grid(column=0, row=2)

        # Title of Source Label
        new_source_title_label = tk.Label(new_source_frame, text="Title of Source:", font=("Arial", 12), width=50,bg="#6EB3F4")
        new_source_title_label.grid(column=0, row=0, pady=5)

        # Title of Source Entry
        self.new_source_title_entry = tk.Entry(new_source_frame)
        self.new_source_title_entry.grid(column=1, row=0, pady=5)

        # Country of Source Label
        new_source_country_label = tk.Label(new_source_frame, text="Country of Source:", font=("Arial", 12), width=50,bg="#6EB3F4")
        new_source_country_label.grid(column=0, row=1, pady=5)

        # Country of Source Entry
        self.new_source_country_entry = tk.Entry(new_source_frame)
        self.new_source_country_entry.grid(column=1, row=1, pady=5)

        # Source Type Label
        source_type_label = tk.Label(new_source_frame, text="Source Type:", font=("Arial", 12),bg="#6EB3F4")
        source_type_label.grid(column=0, row=2, pady=5)

        # Radio button for URL
        self.source_type_url = tk.Radiobutton(new_source_frame, text="URL", bg="#6EB3F4", variable=self.source_type_var, value="url", command=self.toggle_source_input) 
        self.source_type_url.grid(column=1, row=2, pady=5)

        # Radio button for File
        self.source_type_file = tk.Radiobutton(new_source_frame, text="File", bg="#6EB3F4", variable=self.source_type_var, value="file", command=self.toggle_source_input)
        self.source_type_file.grid(column=1, row=3, pady=5)

        # Initially, set the source type to URL
        self.source_type_var.set("url")

        # URL of Source Label
        new_source_url_label = tk.Label(new_source_frame, text="URL:", font=("Arial", 12), bg="#6EB3F4")
        new_source_url_label.grid(column=0, row=4, pady=5)

        # # URL of Source Entry
        self.new_source_url_entry = tk.Entry(new_source_frame)
        self.new_source_url_entry.grid(column=1, row=4, pady=5)

        # # File Upload Button
        self.file_upload_button = tk.Button(new_source_frame, text="Select File", bg="#6EB3F4", command=self.upload_file)
        # self.file_upload_button.grid(column=0, row=4, pady=5)

        # (Endnote classification) Originator; Source identifier; Date of publication; (skipping 1/17) Date of information [optional but preferred]; (Classification of title/subject) Title/Subject; (skipping 1/17) Page/paragraph or portion indicator [required when applicable]; Classification of extracted information is “X”; Overall classification is “X”; Access date.

        # URL of Source Label
        originator_label = tk.Label(new_source_frame, text="Originator:", font=("Arial", 12), bg="#6EB3F4")
        originator_label.grid(column=0, row=5, pady=5)

        # Originator Entry
        self.originator_entry = tk.Entry(new_source_frame)
        self.originator_entry.grid(column=1, row=5, pady=5)

        # Date of Publication Label
        date_publication_label = tk.Label(new_source_frame, text="Date of Publication:", font=("Arial", 12), bg="#6EB3F4")
        date_publication_label.grid(column=0, row=6, pady=5)

        # Date of Publication Entry
        self.date_publication_entry = tk.Entry(new_source_frame)
        self.date_publication_entry.grid(column=1, row=6, pady=5)

        classification_options = [
            "Unclassified",
            "Confidential",
            "Secret",
            "Top Secret",
        ]

        # Class. Title Label
        classification_title_label = tk.Label(new_source_frame, text="Classification of Subject/Title:", font=("Arial", 12),bg="#6EB3F4")
        classification_title_label.grid(column=0, row=7, pady=5)

        # Class. Title Entry
        self.classification_title_var = tk.StringVar()
        self.classification_title = tk.OptionMenu(new_source_frame, self.classification_title_var, *classification_options)
        self.classification_title.grid(column=1, row=7, pady=5)

        # Class. Title Label
        portion_classification_label = tk.Label(new_source_frame, text="Classification of Portion:", font=("Arial", 12),bg="#6EB3F4")
        portion_classification_label.grid(column=0, row=8, pady=5)

        # Class. Portion Entry
        self.portion_classification_var = tk.StringVar()
        self.portion_classification = tk.OptionMenu(new_source_frame, self.portion_classification_var, *classification_options)
        self.portion_classification.grid(column=1, row=8, pady=5)

        # Class. Title Label
        overall_classification_label = tk.Label(new_source_frame, text="Overall Classification:", font=("Arial", 12),bg="#6EB3F4")
        overall_classification_label.grid(column=0, row=9, pady=5)

        # Class. Overall Entry
        self.overall_classification_var = tk.StringVar()
        self.overall_classification = tk.OptionMenu(new_source_frame, self.overall_classification_var, *classification_options)
        self.overall_classification.grid(column=1, row=9, pady=5)


        # Submit Button
        new_source_submit = tk.Button(new_source_frame, text="Add Source", bg="#6EB3F4", command=self.add_new_source_command)
        new_source_submit.grid(column=0, row=10, pady=5)


        #Source List Label
        self.source_list_label = tk.Label(self.root, text="Source List", font=("Arial", 18), bg="#6EB3F4", width=50)
        self.source_list_label.grid(column=1, row=1)

        # Source List Frame
        self.source_list_frame = tk.Frame(self.root,bg="#6EB3F4")
        self.source_list_frame.grid(column=1, row=2, sticky="nsew")



        # Save Summaries Frame 
        self.save_source_summaries_frame = tk.Frame(self.root, bg="#6EB3F4")
        self.save_source_summaries_frame.grid(column=0, row=4)

        # Summarize Button
        self.summarize_button_text_variable = tk.StringVar()
        self.summarize_button_text_variable.set("Summarize Sources")
        self.summarize_button = tk.Button(self.root, textvariable=self.summarize_button_text_variable, bg="#6EB3F4", command=self.summarize)
        self.summarize_button.grid(column=0, row=3, pady=6)
        #Preview Report
        self.preview_report_button = tk.Button(self.save_source_summaries_frame, text="Preview Report", bg="#6EB3F4", command=self.popup)
        self.preview_report_button.grid(column=0, row=1, pady=6)

        # # Use ATS Button
        self.use_ats_button = tk.Button(self.save_source_summaries_frame, text = "Preview ATS", command=self.ats_popup)
        self.use_ats_button.grid(column=1, row=2, pady=10)

        # Save ATS to docx Button
        self.save_ats_to_docx_button =tk.Button(self.save_source_summaries_frame, text = "Save ATS to Word Document", bg="#6EB3F4", command=self.save_ats_to_docx)
        self.save_ats_to_docx_button.grid(column=1, row=3, pady=8)
        # # New Source Label
        # self.save_source_summaries_label = tk.Label(self.save_source_summaries_frame, text="Save Source", font=("Arial", 16),bg="#6EB3F4")
        # self.save_source_summaries_label.grid(column=0, row=0)

        #Save to DOC
        
        self.save_button_text_variable = tk.StringVar()
        self.summarize_button_text_variable.set("Summarize Sources")
        self.save_source_summaries_button = tk.Button(self.save_source_summaries_frame, text="Save Report to Word Document", bg="#6EB3F4", command=self.save_summaries_to_docx)
        self.save_source_summaries_button.grid(column=0, row=2, pady=5)

        #Save to PDF
        self.save_source_summaries_button_pdf = tk.Button(self.save_source_summaries_frame, text="Save Report to to PDF",bg="#6EB3F4", command=self.save_summaries_to_pdf)
        self.save_source_summaries_button_pdf.grid(column=0, row=3, pady=6)

         #Save to pptx
        
        # self.save_pptx_button_text_variable = tk.StringVar()
        # self.summarize_button_text_variable.set("Summarize Sources")
        # self.save_source_summaries_button = tk.Button(self.save_source_summaries_frame, text="Save to PowerPoint", command=self.save_summaries_to_pptx)
        # self.save_source_summaries_button.grid(column=0, row=3, pady=9)


    


         # Disclaimer Label
        self.disclaimer_label = tk.Label(self.root, text="     * Please note that this program uses ChatGPT, and as such no classified data should be input through the system. * \n \n"
                                                                  "    * The analytical standards and writing style used are up to date as of January 2022 but will not update until the system itself is updated.*" , font=("Arial", 14), bg="#6EB3F4") 
        self.disclaimer_label.grid(column=0, row=10, pady=15)




    def toggle_source_input(self):
        source_type = self.source_type_var.get()

        # Hide all widgets initially
        self.new_source_url_entry.grid_forget()
        self.file_upload_button.grid_forget()

        # Show the relevant widget based on the source type
        if source_type == "url":
            self.new_source_url_entry.grid(column=1, row=4, pady=5, padx=5)
        elif source_type == "file":
            self.file_upload_button.grid(column=1, row=4, pady=5, padx=5)

    def create_file_upload_button(self):
        # Create a file upload button
        self.file_upload_button = tk.Button(self.root, text="Select File", bg="#6EB3F4",command=self.upload_file)
        self.file_upload_button.grid(row=4, pady=5, padx=5)

    def upload_file(self):
        # Open a file dialog to select a file
        self.current_source_location = filedialog.askopenfilename(filetypes=[("PDFs", ".pdf"), ("Word documents", ".docx")])


    def reset_source_inputs(self) -> None:
        self.new_source_title_entry.delete(0, 'end')
        self.new_source_url_entry.delete(0, 'end')

    def delete_source_command(self, index: int) -> None:
        # Delete the source at the specified index
        del self.source_list[index]
        
        # Update the GUI to reflect the changes
        self.update_source_list_gui()

    def update_source_list_gui(self) -> None:
        # Clear all widgets in source_list_frame
        for widget in self.source_list_frame.winfo_children():
            widget.destroy()

        # Re-populate source_list_frame with updated source_list
        for i, source in enumerate(self.source_list):
            source_title = source["source_title"]
            
            # Source List Item Title
            source_list_title = tk.Label(self.source_list_frame, text=source_title, font=("Arial", 12),bg="#6EB3F4")
            source_list_title.grid(column=0, row=i, sticky="w")

            # Delete Button
            delete_button = tk.Button(self.source_list_frame, text="Delete", command=lambda i=i: self.delete_source_command(i))
            delete_button.grid(column=1, row=i, padx=5, sticky="e")

    def add_new_source_command(self) -> None:
        title = self.new_source_title_entry.get()

        # self.add_new_source_command["state"]=tk.NORMAL
        # add text variable?


        if self.source_type_var.get() == "url":
            self.current_source_location = self.new_source_url_entry.get()
        else: 
            self.current_source_location = self.current_source_location

        self.source_list.append({
            "source_title": title,
            "source_classification": "UNCLASSIFIED",
            "source_country": self.new_source_country_entry.get(),
            "source_type": self.source_type_var.get(),
            "source_location": self.current_source_location,
            "source_originator": self.originator_entry.get(),
            "source_date_of_publication": self.date_publication_entry.get(),
            "source_classification": self.classification_title_var.get(),
            "source_portion_classification": self.portion_classification_var.get(),
            "source_overall_classification": self.overall_classification_var.get()
        })

        # self.add_new_source_command.set("Add Source")
        # self.add_new_source_command["state"]=tk.NORMAL


        self.update_source_list_gui()
        self.reset_source_inputs()

        # threading.Timer(0, loading_page, args=(self, self.add_new_source_command)).start()


    def get_text_from_url(self, url):
        try:
            # Fetch the HTML content of the webpage
            response = requests.get(url)
            
            # Check if the request was successful (status code 200)
            if response.status_code == 200:
                # Parse the HTML content with BeautifulSoup
                soup = BeautifulSoup(response.text, 'html.parser')
                
                # Extract all the text from the HTML
                all_text = soup.get_text(separator='\n', strip=True).replace("\n", " ").replace("\t", " ")
                
                return all_text
            else:
                # If the request was not successful, print an error message
                
                raise Exception(f"Error: Unable to fetch content from {url}. Status code: {response.status_code}")
                                        
        except Exception as e:
            print(f"Error: Unable to extract text from {url}. {str(e)}")
            raise e
            

    def get_text_from_pdf(self, filename):
        try:
            # Open the PDF file
            pdf_document = fitz.open(filename)
            
            # Initialize an empty string to store the extracted text
            all_text = ""
            
            # Iterate through each page of the PDF
            for page_number in range(pdf_document.page_count):
                # Get the text of the page
                page = pdf_document[page_number]
                text = page.get_text("text").replace("\n", " ").replace("\t", " ")
                
                # Append the text to the result string
                all_text += text + '\n'
            
            # Close the PDF document
            pdf_document.close()
            
            return all_text.strip()  # Remove leading and trailing whitespaces
        except Exception as e:
            print(f"Error: Unable to extract text from {filename}. {str(e)}")
            raise e
        
    def get_text_from_doc(self,filename):
        #check if the file exists
        if not os.path.isfile(filename):
            print(f"Error: File {filename} does not exist.")
            return
        try:
            #Open the Word Doc
            print(filename)
            word_doc = Document(filename)
            all_text = ""

            for para in word_doc.paragraphs:
                text = para.text

                all_text += text + '\n'

            return all_text.strip() #removes all leading and trailing whitespaces
        except Exception as e:
            print(e)
            print(f"Error: Unable to extract text from {filename}. {str(e)}")
            raise e 
        
    def summarize_all(self):
        summaries = []
        for source in self.source_list:
            summary = self.summarize(source)
            summaries.append(summary)
        final_summary = self.combine_summaries(summaries)
        return final_summary

    def summarize(self):
        self.summarize_button_text_variable.set("Thinking...")
        self.summarize_button["state"]=tk.DISABLED
        for source in self.source_list:
            source_type = source["source_type"]
            text = ""

            if source_type == "url":
                try:
                    text = self.get_text_from_url(source["source_location"])
                except Exception as e:
                    print(e)
                    continue
            else:
                try:
                    if ".docx" in source["source_location"]:
                        text = self.get_text_from_doc(source["source_location"])
                    else:
                        text = self.get_text_from_pdf(source["source_location"])

                except Exception as e:
                    print(e)
                    continue

            # Split the text into chunks based on the maximum token length
            max_token_length = 12000

            encoding = self.encoder.encode(text)

            chunks = [encoding[i:i + max_token_length] for i in range(0, len(encoding), max_token_length)]

            # Initialize summary bits
            summary_bits = []

            # Summarize each chunk and print the results
            for i, chunk in enumerate(chunks):
                print(f"Currently processing chunk {i+1}/{len(chunks)}...")
                
                text_chunk = self.encoder.decode(chunk)
                summary_object = self.client.chat.completions.create(
                    messages=[
                        {
                            "role": "system",
                            "content": "You are going to act as a summarizer for the following text, giving 2-3 sentences of summarization:"
                        },
                        {
                            "role": "user",
                            "content": text_chunk
                        }
                    ],
                    model="gpt-3.5-turbo-16k"
                )

                summary_bits.append(summary_object.choices[0].message.content.replace("\n", " "))

                all_summaries_together_text = " ".join(summary_bits)

            #ensures ICD 203 standards

            total_summary_object = self.client.chat.completions.create(
                messages=[
                    {
                        "role": "system",
                        "content": "You are going to act as a summarizer for the following text, giving 1-2 sentences of summarization which encapsulate the key findings of the text. This will be labelled as the BLUF:, followed by 2-4 sentences with more detail. Use relevant intelligence community directives to analyze the text:"

                    },
                    {
                        "role": "user",
                        "content": all_summaries_together_text
                    }
                ],
                model="gpt-3.5-turbo-16k"
            )

            summary_object_analytical_standards = self.client.chat.completions.create(
                messages=[
                    {
                        "role": "system",
                        "content": """
                        You are an intelligence analyst. Utilize the following analytic standards to review the summary reporting object. Analytical Evaluation Prompt: Understanding and Applying Analytic Standards in Intelligence Community Reporting


                            1. **Background:**
                            - Explore the foundational legal documents and directives shaping Analytic Standards, such as the National Security Act of 1947, Intelligence Reform and Terrorism Prevention Act of 2004, Executive Order 12333, and Presidential Policy Directive/PPD-28.


                            2. **Purpose and Scope:**
                            - Summarize the Intelligence Community Directive (ICD) regarding Analytic Standards, outlining its purpose and specifying its applicability to the Intelligence Community (IC) and elements designated by the President or Director of National Intelligence (DNI).


                            3. **Policy Framework:**
                            - Outline the policy framework established by the IC Analytic Standards, highlighting their core principles in intelligence analysis. Discuss the responsibility of IC elements in applying these standards appropriately, considering factors like purpose, source information, production timeline, and customer needs.


                            4. **Analytic Tradecraft Standards:**
                            - Detail the nine Analytic Tradecraft Standards, providing foundational assessment criteria for evaluating IC analytic products. Emphasize their significance in maintaining rigor, excellence, and integrity in analytic practice.


                            5. **Responsibilities:**
                            - Define the responsibilities of key entities, including the Deputy DNI for Intelligence Integration (DDNI/II) and the Chief, Analytic Integrity and Standards Group (ODNI Analytic Ombuds). Highlight the DDNI/II's role in directing the application of Analytic Standards, overseeing periodic reviews, and refining the evaluation program.


                            6. **Confidentiality and Conflict Resolution:**
                            - Explain the role of the ODNI Analytic Ombuds in addressing concerns raised by analysts about adherence to standards. Emphasize the confidentiality of the process, allowing analysts to voice concerns without fear of reprisal, with exceptions for significant misconduct or legal violations.


                            7. **Implementation by IC Elements:**
                            - Discuss the responsibilities of Heads of IC elements in ensuring the proper application of Analytic Standards. Highlight the designation of individuals or offices within IC elements to respond to analyst concerns, emphasizing the importance of internal reviews and evaluations.


                            8. **Training and Education:**
                            - Emphasize the role of Analytic Standards in guiding education and training programs within IC elements. Stress the need for ongoing improvement based on lessons learned from analytic product evaluations.


                            9. **Assessment Criteria:**
                            - Establish criteria for assessing the effectiveness of Analytic Standards in maintaining objectivity, avoiding bias, ensuring timeliness, incorporating all available sources, and exhibiting clear and logical argumentation.


                            10. **Recommendations for Improvement:**
                                - Propose measures for improving the application of Analytic Standards, addressing any identified challenges or gaps in adherence. Consider potential enhancements to education and training programs and the resolution of concerns raised through the Ombuds process.
                                """
                    },
                        
                    {
                        "role": "user",
                        "content": total_summary_object.choices[0].message.content
                    }
                ],
                model="gpt-3.5-turbo-16k"
            )

            summmary_object_formatting = self.client.chat.completions.create(
                messages=[
                    {
                        "role": "system",
                        "content": """
                        You are a military analyst. Internalize the following Morning Intelligence Update (MIU) format in the order provided. You must not write anything until I prompt you.

                    Here's the format: 

                        Unclassified

                        (U): COUNTRY | APLE | Virtual Analyst | DD Month YYYY

                        Notes: “(U)” stands for unclassified, “COUNTRY” should be substituted in each MIU based on the topic of the text you read for instance ISRAEL etc: you do not need the word COUNTRY in the actual header., “APLE | Virtual Analyst” stays the same regardless of MIU topic, “ DD Month YYYY” stands for current day, month, and year ex: 01 January 2024

                        Below this is a paragraph of no more than seven sentences: 

                        (U) The first sentence or “BLUF” aka Bottom Line Up Front should be written in bold, should capture the main issue or development that is of interest to the CENTCOM Commander (the “what”); Write the BLUF in this style: “On DD MON, event or action happened, according to news agency.” Pay attention that the “On DD MON” is the event’s date, not today’s date. If the day is not provided, provide the month, if the month is not provided at least provide the year. 
                        Summarize the text in seven sentences 
                        Sentences must be less than 21 words. 
                        MIU Format should contain no headers except for “(U): COUNTRY | APLE | Virtual Analyst | DD Month YYYY”
                        Only the first sentence, the BLUF, should be bolded. 
                        Important: sentences should follow right after another and not have spaces in between. 
                        The most important thing of this whole MIU is each sentence must follow the other and be connected into one large paragraph. 
                        The second most important thing is that only the second sentence is bolded and in emphasis. 
                        should not be written in bold or emphasis. Sentences cannot be more than 20 words. You may only write six sentences after the BLUF.
                        The MIU Content Paragraph must be one, connected paragraph, no spaces, beginning with the BLUF in bold and ending with an assessment explaining why the development or event matters. 

                        The sentences must elaborate and develop the main issue of an event, and must be drawn primarily from reporting of the main issue of an event. Do not summarize the entire source since you only have six sentences after the BLUF to use. Do not overuse adjectives and adverbs. The final sentence must contain an assessment explaining why the development matters to the CENTCOM Commander (the “so What” aka the main issue).
                        The MIU Content Paragraph must be one, connected paragraph, no spaces, beginning with the BLUF in bold and ending with the assessment explaining why the development or event matters. 

                        Note: need the heading / last sentence should be part of the paragraph need as one whole thing 
                    Find an example utilizing DIA standard as a reference.
                        """

                    },
                    {
                        "role": "user",
                        "content": summary_object_analytical_standards.choices[0].message.content
                    }
                ],
                model="gpt-3.5-turbo-16k"
            )

            summary_object_JAC = self.client.chat.completions.create(
                messages=[
                    {
                        "role": "system",
                        "content": """

                        Review the summary object generated, and analyze it against this checklist. 
                        Title:
                    [Provide a clear and concise title that reflects the subject of the intelligence product]

                    Objective:
                    Clearly state the primary objective of the intelligence product. What specific information or analysis is the product intended to convey? Ensure that the objective is narrowly defined and aligned with the purpose of the analysis.

                    Political Consideration:
                    Confirm that the analysis is independent of political considerations. Explicitly state any potential political implications or biases and ensure the content remains objective, unbiased, and focused on factual analysis.

                    Timeliness:
                    Verify the timeliness of the intelligence product. Assess the relevance of the information within the current geopolitical or operational context. Provide a timeline for key events or developments that contribute to the timeliness of the analysis.

                    Sources:
                    Describe in detail the quality and credibility of underlying sources, data, and methodologies used in the analysis. Include not only basic, generic descriptions of cited reporting but also provide insights into the methodologies employed. Clearly identify which sources are deemed most important to major analytic judgments.

                    Uncertainties:
                    Express and explain uncertainties associated with major analytic judgments. Indicate levels of uncertainty and explain their basis. Provide a thorough discussion of the nature and sources of uncertainties affecting major analytic judgments. Identify indicators that, if detected, would alter levels of uncertainty.

                    Distinctions:
                    Ensure consistent distinctions among statements conveying underlying information, assumptions, and judgments. Explicitly state linchpin assumptions that serve as the foundation for the argument. Identify indicators that, if detected, could validate or refute judgments or assumptions. Clearly explain the implications for judgments if assumptions are proven incorrect.

                    Alternatives:
                    Incorporate a detailed analysis of alternatives. Present alternatives when uncertainties, complexity, or low probability/high impact situations warrant inclusion. Explain the evidence and reasoning behind each alternative. Discuss the likelihood or implications of each alternative related to U.S. interests. Identify indicators that, if detected, would affect the likelihood of any identified alternatives.

                    Argumentation:
                    Ensure the intelligence product uses clear and logical argumentation. Clearly present the main analytic message, ensuring it is prominent and aligned with the objective. Verify that reasoning is flawless and effectively combines evidence, context, and assumptions to support analytic judgments. Use clear language and a structure that displays a logical flow appropriate for the argument being presented.

                    Change or Consistency:
                    Explain any change to or consistency of analytic judgments. Clearly note how a major analytic judgment compares with previous production and explain how new information or reasoning supports changing or maintaining an existing analytic line. Highlight and explain how a major analytic judgment compares with judgments on the topic within the U.S. intelligence community, not just within the same analytic element.

                    Accuracy:
                    Verify that judgments or assessments are expressed clearly and conditioned. Ensure that each judgment or assessment is clearly articulated and qualified, using "if/then" statements when necessary.

                    Visual Information:
                    Incorporate effective visual information where appropriate. Ensure visual elements are not only pertinent but also clarify, complement, or augment data or analytic points in an effective manner. Take particularly effective advantage of visual presentational methods to convey data or analytic points in a way that enhances the product’s value by making complex issues more understandable, adding insight or perspective, increasing knowledge retention, or highlighting trends, drivers, or indicators.

                    Review and Refinement:
                    Review the generated intelligence product against each element in the checklist. Refine the content as needed to meet the specified standards. Consider feedback from team leads and subject matter experts in the refinement process.

                        """

                    },
                    {
                        "role": "user",
                        "content": summmary_object_formatting.choices[0].message.content
                    }
                ],
                model="gpt-3.5-turbo-16k"
            )
            
            self.formatted_total_summary = summmary_object_formatting.choices[0].message.content


                # Create Citation
            source_overall_classification = source['source_overall_classification']
            source_cit_let = source_overall_classification[0] if source_overall_classification else ''

            current_datetime = datetime.datetime.now()
            date_accessed = current_datetime.strftime("%d/%m/%Y")
            source_citation = f"({source_cit_let});{source['source_originator']}; {source['source_country']}; {source['source_date_of_publication']}; ({source['source_classification']}) {source['source_title']}; Classification of extracted information is {source['source_portion_classification']}; Overall classification: {source['source_overall_classification']}, {date_accessed}"

            summary_dict_object = {
                "source_title": source["source_title"],
                "source_classification": source["source_classification"],
                "source_country": source["source_country"],
                "source_type": source["source_type"],
                "source_summary": summary_object_JAC.choices[0].message.content,
                "source_citation": source_citation
            }

            self.summaries.append(summary_dict_object)
    # #    # Assuming 'self.summaries' is a list containing summary objects
    #     for summary in self.summaries:
    #             print(summary)


        print(f"Done Summarizing Source: {source['source_title']}")

        print("Done summarizing all sources")
        self.summarize_button_text_variable.set("Sources Summarized")
        self.summarize_button["state"]=tk.NORMAL

# creates window to view report in. 
    def popup(self):
        popupwindow = tk.Toplevel(self.window)
        popupwindow.title("Virtual Analyst Preview Window")
        popupwindow.geometry("200x100")
        report = ["source_summary" "42"]
        for summary in self.summaries: print(summary)
        Preview = tk.Label(popupwindow, text = report["source_summary"])
        button1 = tk.Button(popupwindow, text = "Done", command=popupwindow.destroy)
        Preview.pack()
        button1.pack()


# Creates popup to view ATS

    def ats_popup(self):
            ats_popup_window = tk.Toplevel(self.window)
            ats_popup_window.title("Analytic Tradecraft Standards for Virtual Report")
            ats_popup_window.geometry("200x100")

            # PROBLEM: We need some data from the summary step to implement and build an ATS. 
            # EITHER: We save the entire ATS in the summary step (wasteful, since there's a possibility it won't be used)
            # OR: We save the formatted total summary to be used in this step <- we're trying this
            if self.formatted_total_summary is None:
                return
            
            ats_object = self.client.chat.completions.create(
                messages=[
                    {
                        "role": "system",
                        "content": """Do not respond. internalize these standards and format to analyze a text provided separately
                        (U) Analytic Tradecraft Summary

                        (Class) Confidence Level: 

                        State your confidence level in the main assessment 

                        explain the basis for it by referencing your sources of uncertainty including strengths and weaknesses in the information base, assumptions, gaps, alternatives, and the complexity of the issue. Use the following guide to determine your confidence level, and choose the confidence level that best reflects the text. 

                        Signs that indicate high confidence: 
                        Well-corroborated information from proven sources; minimal contradictory reporting; low potential for deception; few information gaps. 
                        Assumptions would not have a significant effect on the assessment if incorrect. 
                        Very unlikely alternative. 
                        Routine event that is well understood; relatively few variables. 

                        Signs that indicate medium confidence:
                        Partially corroborated information from good sources; some potential for deception; several gaps in the information base. 
                        Plausible, yet unlikely, alternatives. 
                        Key assumptions with potentially substantial effect on the assessment if incorrect. 
                        More complex situation with multiple issues or actors; some previous examples that are well understood. 

                        Signs that indicate low confidence:
                        Uncorroborated information; high potential for deception; many critical gaps in the information base. 
                        Plausible alternatives with a nearly even chance of occurring. 
                        Key assumptions with substantial effect on the assessment if incorrect. 
                        Highly complex or rapidly evolving situation with multiple issues or actors; few previous examples that are not well understood.
















                        (CLASS) Sourcing: Provide considerable detail on the strengths and weaknesses of the reporting used for the assessment, focusing on the credibility and quality of the sourcing. Do not provide a summary of what each intelligence type (INT) told you for the analysis. Do not provide a summary of what the reporting said. Identify which sources/reporting was the most important to the assessment and judgment.

                        Elements of source characterization:
                        Context 
                        When did the reported information occur? 
                        What are the source’s strengths or weaknesses? (subject matter expertise, biases, possible denial and deception, etc.) 
                        Credibility/Quality 
                        Is the information credible? Is it of good quality? (accurate, consistent with other reporting, plausible given circumstances) 
                        Reliability 
                        Is the source reliable? (vetted, history of reliable reporting) 
                        Access 
                        How close is the source to the information? (first-, secondhand, further removed 
                        Source Types 
                        Who told us? (informant) 
                        Who did/said? (actor) 
                        What is the origin of the reporting? (type of intelligence) 

                        (CLASS) Gaps: Include intelligence gaps along with a description about the extent to which filling that intelligence gap would alter or bolster your assessment. Intelligence gaps must be tied to your main assessment if those gaps are critical to and underpin the main judgment. 

                        Characterize the extent and limits of your knowledge base. What are some remaining intelligence gaps that prevent you from making a stronger or more useful assessment and that are not explicitly covered by an assumption or judgment?

                        (CLASS) Assumptions: Clearly state a linchpin or gap-bridging assumption(s) underpinning the main assessment. An assumption should help readers understand the connecting tissue between the evidence and the assessment; it generally is something that an analyst believes to be true, but lacks evidence, and if incorrect, would force a change to the assessment. When crafting an assumption, think along the lines of “what would change my assessment?” Identify indicators that could validate or refute assumptions and explain the implication for judgements in assumptions that are incorrect.

                        Internalize the different assumptions of the text and identify which ones are apparent. 

                        Framing assumptions:
                        What longstanding analytic lines are related to the assessment?

                        What beliefs do we hold about what “will always,” “will never,” or
                        “generally will” occur, or what “has always” or “has never” been done
                        or happened, relative to the intelligence question?

                        Do we have a default mindset in how we approach this problem?

                        What are the relevant historical precedents for this question?

                        Have we identified any trends that we expect to continue?

                        Scoping assumptions:
                        What factors, drivers, or variables are not included in the analysis?

                        What factors are we “holding constant” and assuming will
                        not change?

                        Have we assumed that certain events will or will not take place
                        or that certain factors will or will not change?

                        Have we clarified which actors, events, and timeframes are and are
                        not included?

                        Have we defined all of the key terms and concepts in our analysis?

                        Evidence assumptions:
                        Are there multiple possible interpretations of our evidence?

                        Why do we lean toward one interpretation rather than another?

                        What beliefs do we hold about our information base that lead us
                        to ascribe more value to certain pieces of information?

                        What are our beliefs about the extent of our access to all
                        relevant information?

                        Logic assumptions:
                        Have we used a small sample to infer something about
                        a broader group?

                        Have we extrapolated from a known situation to an unknown
                        situation?

                        Do we believe that certain types of events or activities are
                        symptomatic of or more/less likely to occur with some wider
                        phenomenon or conclusion?

                        Do we believe that one event or factor is causing or affecting
                        another event?

                        Bridging assumptions:
                        What are the essential elements of information needed to answer the
                        intelligence question? For which elements do we lack evidence? For
                        which do we have significant uncertainty?

                        What are the factors or conditions that must be present for the
                        assessment to be true (or false), and do we have evidence that they
                        are (or are not) present?

                        For each assumption type, determine if they are high impact assumptions, low impact assumptions, and also determine if they are assumptions that are weak or strong. Internalize definitions below. 

                        (U) High-impact assumptions, if proved false, invalidate or significantly alter
                        the assessment.

                        (U) Low-impact assumptions, if proved false, change only an aspect of the assessment,
                        such as the scope, specificity, likelihood, or timeframe.

                        (U) An assumption is weak or vulnerable if we can imagine a plausible situation, or
                        multiple situations, in which the assumption might not be true.

                        (U) An assumption is strong if we have difficulty imagining a situation in which
                        the assumption might not be true because such a circumstance is highly unlikely
                        or implausible.

                        Based on what you read, choose between high-impact assumptions and low-impact assumptions to characterize the text. Also choose between if the assumption is either weak or strong. 

                        (CLASS) Alternatives: Include a plausible and useful alternative to your main assessment. Explain the reasoning and/or evidence that underpins the alternatives. Discuss the alternative likelihood or implications related to United States interests. Identify indicators that, if identified, would affect the likelihood of the alternatives.

                        First, identify the sources of the uncertainty that bound our understanding of a problem set. This
                        can be done as part of a regularly occurring analytic line review or in support of a specific project.
                        The following questions can identify gaps, assumptions, or different interpretations of evidence that
                        can generate alternatives: 

                        What prevents us from being absolutely certain that our authoritative judgment is correct?

                        What limits our confidence level?

                        What are the assessment’s underlying assumptions, and under what conditions might
                        they prove false?

                        What are the weaknesses of our information base?

                        Is there any contradictory reporting?

                        Could there be denial and deception, deliberate falsification, or other misinformation
                        that could affect our analysis?

                        Are we over relying, or relying exclusively, on one intelligence collection stream
                        or platform?

                        Second, consider alternative hypotheses. Discussing these can enable better detection of future
                        events or developments that change the authoritative assessment. Ask:

                        What other hypotheses or options did we consider, and are they plausible?

                        Are there multiple explanations for the information we have?

                        How vulnerable is the assessment to change?

                        What would have to change to make us reconsider the expected outcome?

                        What indicators of change would we expect to be captured with our collection assets?

                        What indicators do we think we could not observe?

                        Finally, consider the implications of our assessments for our clients in order to mitigate surprise,
                        allow for planning, and provide warning:

                        What are the implications for U.S. interests if we are wrong about our assessment?
                        •
                        What types of plausible events would be game changers, that is, would fundamentally shift
                        the issues of import or outcomes we currently anticipate? What would the implications be?

                        How would we know that our alternative is becoming likely or that our authoritative assessment
                        is becoming unlikely?

                        Next address the elements below to ensure the alternatives presented in every product are useful, plausible, and rigorous.

                        What is the alternative to the authoritative assessment?

                        What is its likelihood (relative and absolute)? Some alternatives may be highly unlikely, whereas
                        others may not be significantly less likely than the authoritative assessment.

                        What reasoning and/or evidence substantiates the plausibility of the alternative? Explain
                        the support for the alternative, rather than using the alternative to bolster the case for the
                        authoritative assessment.

                        What are the implications for U.S. interests of the alternative that warrant consideration?

                        When appropriate, what indicators would, if observed, affect the likelihood of the alternative
                        and the authoritative assessment?

                        Internalize approaches to writing alternative assessments. 

                        Exploring the Potential for Surprise. This approach to analysis of alternatives examines the
                        impact of a hard-to-predict event or a surprise to facilitate contingency planning. It includes
                        collectible, specific indicators to providing warning.

                        Competing Assessments. This approach clarifies the alternative’s strengths and weaknesses
                        as compared with the authoritative assessment. This type of alternative can be, but is not limited
                        to, a competing view from another IC element. Addressing alternatives can enhance the credibility
                        of our assessments.

                        Discussing the Implications of Information or Assumptions. This approach examines the
                        impact of key information or assumptions on our judgments, allowing clients to determine
                        whether contingency planning is needed. Indicators are highlighted, as appropriate, in the product.

                        These directions internalized write at least two alternatives to the initial assessment in the text. 

                        :"""

                    },
                    {
                        "role": "user",
                        "content":  self.formatted_total_summary
                    }
                ],
                model="gpt-3.5-turbo-16k"
            )

            ats = ats_object.choices[0].message.content

            include = tk.Label(ats_popup_window, text = ats, width=60)
            buttoni = tk.Button(ats_popup_window, text = "Done", command=ats_popup_window.destroy)
            include.pack()
            buttoni.pack()

        






# Call loading_page from the main thread
# threading.Timer(0, loading_page, args=(self, self.summarize)).start()
    
#designed to combine multiple reports into one
# def combine_summaries(self, summaries):
#     # Combine all summaries into one final summary
#     final_summary = " ".join(summaries)
#     return final_summary

       # creating pptx code in progress              

    # def save_summaries_to_pptx(self):
    #     self.summarize_button_text_variable.set("Thinking...")
    #     self.summarize_button["state"]=tk.DISABLED

    #     prs = Presentation()

    #     print("Report Saved to PowerPoint File")
    #     title_slide_layout = prs.slide_layouts[0]
    #     slide = prs.slides.add_slide(title_slide_layout)
    #     title = slide.shapes.title
    #     subtitle = slide.placeholders[1]

    #     title.text = "[Intelligence Note or Reporting Highlights]"
    #     subtitle.text = datetime, "CCJ2 APLE:"

        # for summary in self.summaries:
        #     slide_layout = prs.slide_layouts[1]
        #     slide = prs.slides.add_slide(slide_layout)
        #     shapes = slide.shapes
        #     title_shape = shapes.title
        #     content_shape = shapes.placeholders[1]
            

        #     title_shape.text = summary["source_title"]
        #     content_shape.text = summary["source_summary"]
            
        #     textbox = slide.shapes[0]

        #     text_frame = textbox.text_frame
        #     text_frame.add_paragraph(content_shape.text)

        #     endnotes_slide = slide.notes_slide
        #     notes_text_frame = endnotes_slide.notes_text_frame

        #     citation = "Endnotes:"
        #     notes_text_frame.add_paragraph().add_run(citation).font.size = (10)

        # prs.save("ReportTemplate.pptx")
        # self.summarize_button_text_variable.set("PowerPoint Saved")
        # self.summarize_button["state"] = tk.NORMAL



#loading screen code in progress

# # Define root as a Tkinter window
# root = tkinter.Tk()
# # Now that the ReportGeneratorApp class has been defined, we can create an instance of it
# app = ReportGeneratorApp(root)

# the_queue = queue.Queue()

# def worker_thread(app, the_queue):
#    try:
#        # Perform a long-running operation...
#        result = app.summarize()
#        result = app.add_new_source_command
#        result = app.save_summaries_to_docx
#        result = app.save_summaries_to_pdf

#        # Put the result onto the queue
#        the_queue.put(result)
#    except Exception as e:
#        print(f"An error occurred in worker_thread: {e}")

#    # Start the worker thread
# threading.Thread(target=worker_thread, args=(app,)).start()


# def update_interface():
#    # Check if there's anything on the queue
#    try:
#        result = the_queue.get(block=False)
#    except queue.Empty:
#        # No data on the queue, schedule another update
#        root.after(100, update_interface)
#        return

#    # Update the Tkinter interface with the result
#    label.config(text=result)

#    # Schedule another update
#    root.after(100, update_interface)

# # root = tkinter.Tk()
# # root.title("Checking the queue")
# # label = tkinter.Label(root, text="Waiting for result...")
# # label.pack()

# # Start the worker thread
# threading.Thread(target=worker_thread).start()

# # Start updating the interface
# update_interface()


# # Create a single Tkinter window for the loading screen
# root = tk.Tk()
# root.title("Loading Page")
# root.geometry("300x200")
# root.resizable(False, False)

# # Create a label for the loading message
# loading_label = ttk.Label(root, text="Thinking...", font=("Arial", 14))
# loading_label.pack(pady=50)

# # Create a style object
# s = ttk.Style()

# # Configure progress bar style
# s.configure("blue.Horizontal.TProgressbar", foreground='blue', background='blue')

# # Create a progress bar with the custom style
# progress_bar = ttk.Progressbar(root, style="blue.Horizontal.TProgressbar", length=200, mode='indeterminate')
# progress_bar.pack()

# # Hide the window initially
# root.withdraw()


# loading_screen = False

# def loading_page(functions):
#    for func in functions:
#        # Display the loading screen
#        root.deiconify()

#        # Create a new thread to run the function
#        operation_thread = threading.Thread(target=func)
#        operation_thread.start()

#        # Wait for the function to finish
#        operation_thread.join()

#        # Hide the loading screen
#        root.withdraw()



if __name__ == "__main__":
   root = tk.Tk()
   app = ReportGeneratorApp(root)
#    loading_page([app.summarize, app.save_summaries_to_docx, app.save_summaries_to_pdf, app.add_new_source_command])
   root.mainloop()

